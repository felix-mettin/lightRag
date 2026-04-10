"""
This module contains all document-related routes for the LightRAG API.
"""

import asyncio
import re
from functools import lru_cache
from lightrag.utils import logger, get_pinyin_sort_key
import aiofiles
import traceback
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List, Optional, Any, Literal, cast
from io import BytesIO
from fastapi import (
    APIRouter,
    BackgroundTasks,
    Depends,
    File,
    Form,
    HTTPException,
    UploadFile,
    Query,
)
from pydantic import BaseModel, Field, field_validator

from lightrag import LightRAG
from lightrag.base import DeletionResult, DocProcessingStatus, DocStatus
from lightrag.utils import (
    generate_track_id,
    compute_mdhash_id,
    sanitize_text_for_encoding,
)
from lightrag.api.utils_api import get_combined_auth_dependency
from ..config import global_args


# 根据文件名识别标准类型, 返回"GB"、"HB"、"GJB"或"others"

def detect_standard_type(filename: str) -> str:    
    """根据文件名识别标准类型：国标/行标/国际标/others"""
    # 直接使用文件名，不尝试移除路径（因为文件名可能包含 /）
    upper_name = filename.upper()

    # 国标：GB、GB/T、GB T、GB-T、GB_T 后跟数字，以及 GBZ 等
    if re.search(r'GB[\s\-_/.]*T?[\s\-_/.]*\d', upper_name) or \
       re.search(r'^GBZ?[\s\-_]?\d', upper_name) or \
       '国家标准' in upper_name:
        return "GB"

    # 国际标（优先于行标）
    if re.search(r'^IEC[\s\-_]?\d', upper_name) or \
       re.search(r'^ISO[\s\-_]?\d', upper_name) or \
       re.search(r'^IEEE[\s\-_]?\d', upper_name) or \
       '国际标准' in upper_name:
        return "GJB"

    # 行标：JB、DL、YD、HG、SY、QB、FZ 等代号，后跟可选分隔符和可选 T
    if re.search(r'^(JB|DL|YD|HG|SY|QB|FZ)[\s\-_./]?T?[\s\-_./]?\d', upper_name) or \
       '行业标准' in upper_name or '行标' in upper_name:
        return "HB"

    return "others"


@lru_cache(maxsize=1)
def _is_docling_available() -> bool:
    """Check if docling is available (cached check).

    This function uses lru_cache to avoid repeated import attempts.
    The result is cached after the first call.

    Returns:
        bool: True if docling is available, False otherwise
    """
    try:
        import docling  # noqa: F401  # type: ignore[import-not-found]

        return True
    except ImportError:
        return False


# Function to format datetime to ISO format string with timezone information
def format_datetime(dt: Any) -> Optional[str]:
    """Format datetime to ISO format string with timezone information

    Args:
        dt: Datetime object, string, or None

    Returns:
        ISO format string with timezone information, or None if input is None
    """
    if dt is None:
        return None
    if isinstance(dt, str):
        return dt

    # Check if datetime object has timezone information
    if isinstance(dt, datetime):
        # If datetime object has no timezone info (naive datetime), add UTC timezone
        if dt.tzinfo is None:
            dt = dt.replace(tzinfo=timezone.utc)

    # Return ISO format string with timezone information
    return dt.isoformat()


router = APIRouter(
    prefix="/documents",
    tags=["documents"],
)

# Temporary file prefix
temp_prefix = "__tmp__"


def sanitize_filename(filename: str, input_dir: Path) -> str:
    """
    Sanitize uploaded filename to prevent Path Traversal attacks.

    Args:
        filename: The original filename from the upload
        input_dir: The target input directory

    Returns:
        str: Sanitized filename that is safe to use

    Raises:
        HTTPException: If the filename is unsafe or invalid
    """
    # Basic validation
    if not filename or not filename.strip():
        raise HTTPException(status_code=400, detail="Filename cannot be empty")

    # Remove path separators and traversal sequences
    clean_name = filename.replace("/", "").replace("\\", "")
    clean_name = clean_name.replace("..", "")

    # Remove control characters and null bytes
    clean_name = "".join(c for c in clean_name if ord(c) >= 32 and c != "\x7f")

    # Remove leading/trailing whitespace and dots
    clean_name = clean_name.strip().strip(".")

    # Check if anything is left after sanitization
    if not clean_name:
        raise HTTPException(
            status_code=400, detail="Invalid filename after sanitization"
        )

    # Verify the final path stays within the input directory
    try:
        final_path = (input_dir / clean_name).resolve()
        if not final_path.is_relative_to(input_dir.resolve()):
            raise HTTPException(status_code=400, detail="Unsafe filename detected")
    except (OSError, ValueError):
        raise HTTPException(status_code=400, detail="Invalid filename")

    return clean_name


class ScanResponse(BaseModel):
    """Response model for document scanning operation

    Attributes:
        status: Status of the scanning operation
        message: Optional message with additional details
        track_id: Tracking ID for monitoring scanning progress
    """

    status: Literal["scanning_started"] = Field(
        description="Status of the scanning operation"
    )
    message: Optional[str] = Field(
        default=None, description="Additional details about the scanning operation"
    )
    track_id: str = Field(description="Tracking ID for monitoring scanning progress")

    class Config:
        json_schema_extra = {
            "example": {
                "status": "scanning_started",
                "message": "Scanning process has been initiated in the background",
                "track_id": "scan_20250729_170612_abc123",
            }
        }


class ReprocessResponse(BaseModel):
    """Response model for reprocessing failed documents operation

    Attributes:
        status: Status of the reprocessing operation
        message: Message describing the operation result
        track_id: Always empty string. Reprocessed documents retain their original track_id.
    """

    status: Literal["reprocessing_started"] = Field(
        description="Status of the reprocessing operation"
    )
    message: str = Field(description="Human-readable message describing the operation")
    track_id: str = Field(
        default="",
        description="Always empty string. Reprocessed documents retain their original track_id from initial upload.",
    )

    class Config:
        json_schema_extra = {
            "example": {
                "status": "reprocessing_started",
                "message": "Reprocessing of failed documents has been initiated in background",
                "track_id": "",
            }
        }


class CancelPipelineResponse(BaseModel):
    """Response model for pipeline cancellation operation

    Attributes:
        status: Status of the cancellation request
        message: Message describing the operation result
    """

    status: Literal["cancellation_requested", "not_busy"] = Field(
        description="Status of the cancellation request"
    )
    message: str = Field(description="Human-readable message describing the operation")

    class Config:
        json_schema_extra = {
            "example": {
                "status": "cancellation_requested",
                "message": "Pipeline cancellation has been requested. Documents will be marked as FAILED.",
            }
        }


class InsertTextRequest(BaseModel):
    """Request model for inserting a single text document

    Attributes:
        text: The text content to be inserted into the RAG system
        file_source: Source of the text (optional)
    """

    text: str = Field(
        min_length=1,
        description="The text to insert",
    )
    file_source: Optional[str] = Field(default=None, min_length=0, description="File Source")

    @field_validator("text", mode="after")
    @classmethod
    def strip_text_after(cls, text: str) -> str:
        return text.strip()

    @field_validator("file_source", mode="after")
    @classmethod
    def strip_source_after(cls, file_source: Optional[str]) -> Optional[str]:
        return file_source.strip() if file_source is not None else None

    class Config:
        json_schema_extra = {
            "example": {
                "text": "This is a sample text to be inserted into the RAG system.",
                "file_source": "Source of the text (optional)",
            }
        }


class InsertTextsRequest(BaseModel):
    """Request model for inserting multiple text documents

    Attributes:
        texts: List of text contents to be inserted into the RAG system
        file_sources: Sources of the texts (optional)
    """

    texts: list[str] = Field(
        min_length=1,
        description="The texts to insert",
    )
    file_sources: Optional[list[str]] = Field(
        default=None, min_length=0, description="Sources of the texts"
    )

    @field_validator("texts", mode="after")
    @classmethod
    def strip_texts_after(cls, texts: list[str]) -> list[str]:
        return [text.strip() for text in texts]

    @field_validator("file_sources", mode="after")
    @classmethod
    def strip_sources_after(cls, file_sources: Optional[list[str]]) -> Optional[list[str]]:
        if file_sources is None:
            return None
        return [file_source.strip() for file_source in file_sources]

    class Config:
        json_schema_extra = {
            "example": {
                "texts": [
                    "This is the first text to be inserted.",
                    "This is the second text to be inserted.",
                ],
                "file_sources": [
                    "First file source (optional)",
                ],
            }
        }


class InsertResponse(BaseModel):
    """Response model for document insertion operations

    Attributes:
        status: Status of the operation (success, duplicated, partial_success, failure)
        message: Detailed message describing the operation result
        track_id: Tracking ID for monitoring processing status
    """

    status: Literal["success", "duplicated", "partial_success", "failure"] = Field(
        description="Status of the operation"
    )
    message: str = Field(description="Message describing the operation result")
    track_id: str = Field(description="Tracking ID for monitoring processing status")

    class Config:
        json_schema_extra = {
            "example": {
                "status": "success",
                "message": "File 'document.pdf' uploaded successfully. Processing will continue in background.",
                "track_id": "upload_20250729_170612_abc123",
            }
        }


class ClearDocumentsResponse(BaseModel):
    """Response model for document clearing operation

    Attributes:
        status: Status of the clear operation
        message: Detailed message describing the operation result
    """

    status: Literal["success", "partial_success", "busy", "fail"] = Field(
        description="Status of the clear operation"
    )
    message: str = Field(description="Message describing the operation result")

    class Config:
        json_schema_extra = {
            "example": {
                "status": "success",
                "message": "All documents cleared successfully. Deleted 15 files.",
            }
        }


class ClearCacheRequest(BaseModel):
    """Request model for clearing cache

    This model is kept for API compatibility but no longer accepts any parameters.
    All cache will be cleared regardless of the request content.
    """

    class Config:
        json_schema_extra = {"example": {}}


class ClearCacheResponse(BaseModel):
    """Response model for cache clearing operation

    Attributes:
        status: Status of the clear operation
        message: Detailed message describing the operation result
    """

    status: Literal["success", "fail"] = Field(
        description="Status of the clear operation"
    )
    message: str = Field(description="Message describing the operation result")

    class Config:
        json_schema_extra = {
            "example": {
                "status": "success",
                "message": "Successfully cleared cache for modes: ['default', 'naive']",
            }
        }


"""Response model for document status

Attributes:
    id: Document identifier
    content_summary: Summary of document content
    content_length: Length of document content
    status: Current processing status
    created_at: Creation timestamp (ISO format string)
    updated_at: Last update timestamp (ISO format string)
    chunks_count: Number of chunks (optional)
    error: Error message if any (optional)
    metadata: Additional metadata (optional)
    file_path: Path to the document file
"""


class DeleteDocRequest(BaseModel):
    doc_ids: List[str] = Field(..., description="The IDs of the documents to delete.")
    standard_type: Optional[str] = Field(default=None, description="文档所属标准类型，不填则自动识别")
    delete_file: bool = Field(
        default=False,
        description="Whether to delete the corresponding file in the upload directory.",
    )
    delete_llm_cache: bool = Field(
        default=False,
        description="Whether to delete cached LLM extraction results for the documents.",
    )

    @field_validator("doc_ids", mode="after")
    @classmethod
    def validate_doc_ids(cls, doc_ids: List[str]) -> List[str]:
        if not doc_ids:
            raise ValueError("Document IDs list cannot be empty")

        validated_ids = []
        for doc_id in doc_ids:
            if not doc_id or not doc_id.strip():
                raise ValueError("Document ID cannot be empty")
            validated_ids.append(doc_id.strip())

        # Check for duplicates
        if len(validated_ids) != len(set(validated_ids)):
            raise ValueError("Document IDs must be unique")

        return validated_ids


class DeleteEntityRequest(BaseModel):
    entity_name: str = Field(..., description="The name of the entity to delete.")
    standard_type: str = Field(..., description="The workspace/standard type containing the entity.")

    @field_validator("entity_name", mode="after")
    @classmethod
    def validate_entity_name(cls, entity_name: str) -> str:
        if not entity_name or not entity_name.strip():
            raise ValueError("Entity name cannot be empty")
        return entity_name.strip()

    @field_validator("standard_type", mode="after")
    @classmethod
    def validate_standard_type(cls, standard_type: str) -> str:
        if not standard_type or not standard_type.strip():
            raise ValueError("standard_type cannot be empty")
        return standard_type.strip()


class DeleteRelationRequest(BaseModel):
    source_entity: str = Field(..., description="The name of the source entity.")
    target_entity: str = Field(..., description="The name of the target entity.")
    standard_type: str = Field(..., description="The workspace/standard type containing the relation.")

    @field_validator("source_entity", "target_entity", mode="after")
    @classmethod
    def validate_entity_names(cls, entity_name: str) -> str:
        if not entity_name or not entity_name.strip():
            raise ValueError("Entity name cannot be empty")
        return entity_name.strip()

    @field_validator("standard_type", mode="after")
    @classmethod
    def validate_relation_standard_type(cls, standard_type: str) -> str:
        if not standard_type or not standard_type.strip():
            raise ValueError("standard_type cannot be empty")
        return standard_type.strip()


class DocStatusResponse(BaseModel):
    id: str = Field(description="Document identifier")
    content_summary: str = Field(description="Summary of document content")
    content_length: int = Field(description="Length of document content in characters")
    status: DocStatus = Field(description="Current processing status")
    created_at: Optional[str] = Field(description="Creation timestamp (ISO format string)")
    updated_at: Optional[str] = Field(description="Last update timestamp (ISO format string)")
    track_id: Optional[str] = Field(
        default=None, description="Tracking ID for monitoring progress"
    )
    chunks_count: Optional[int] = Field(
        default=None, description="Number of chunks the document was split into"
    )
    error_msg: Optional[str] = Field(
        default=None, description="Error message if processing failed"
    )
    metadata: Optional[dict[str, Any]] = Field(
        default=None, description="Additional metadata about the document"
    )
    file_path: str = Field(description="Path to the document file")

    class Config:
        json_schema_extra = {
            "example": {
                "id": "doc_123456",
                "content_summary": "Research paper on machine learning",
                "content_length": 15240,
                "status": "processed",
                "created_at": "2025-03-31T12:34:56",
                "updated_at": "2025-03-31T12:35:30",
                "track_id": "upload_20250729_170612_abc123",
                "chunks_count": 12,
                "error": None,
                "metadata": {"author": "John Doe", "year": 2025},
                "file_path": "research_paper.pdf",
            }
        }


class DocsStatusesResponse(BaseModel):
    """Response model for document statuses

    Attributes:
        statuses: Dictionary mapping document status to lists of document status responses
    """

    statuses: Dict[DocStatus, List[DocStatusResponse]] = Field(
        default_factory=dict,
        description="Dictionary mapping document status to lists of document status responses",
    )

    class Config:
        json_schema_extra = {
            "example": {
                "statuses": {
                    "PENDING": [
                        {
                            "id": "doc_123",
                            "content_summary": "Pending document",
                            "content_length": 5000,
                            "status": "pending",
                            "created_at": "2025-03-31T10:00:00",
                            "updated_at": "2025-03-31T10:00:00",
                            "track_id": "upload_20250331_100000_abc123",
                            "chunks_count": None,
                            "error": None,
                            "metadata": None,
                            "file_path": "pending_doc.pdf",
                        }
                    ],
                    "PREPROCESSED": [
                        {
                            "id": "doc_789",
                            "content_summary": "Document pending final indexing",
                            "content_length": 7200,
                            "status": "preprocessed",
                            "created_at": "2025-03-31T09:30:00",
                            "updated_at": "2025-03-31T09:35:00",
                            "track_id": "upload_20250331_093000_xyz789",
                            "chunks_count": 10,
                            "error": None,
                            "metadata": None,
                            "file_path": "preprocessed_doc.pdf",
                        }
                    ],
                    "PROCESSED": [
                        {
                            "id": "doc_456",
                            "content_summary": "Processed document",
                            "content_length": 8000,
                            "status": "processed",
                            "created_at": "2025-03-31T09:00:00",
                            "updated_at": "2025-03-31T09:05:00",
                            "track_id": "insert_20250331_090000_def456",
                            "chunks_count": 8,
                            "error": None,
                            "metadata": {"author": "John Doe"},
                            "file_path": "processed_doc.pdf",
                        }
                    ],
                }
            }
        }


class TrackStatusResponse(BaseModel):
    """Response model for tracking document processing status by track_id

    Attributes:
        track_id: The tracking ID
        documents: List of documents associated with this track_id
        total_count: Total number of documents for this track_id
        status_summary: Count of documents by status
    """

    track_id: str = Field(description="The tracking ID")
    documents: List[DocStatusResponse] = Field(
        description="List of documents associated with this track_id"
    )
    total_count: int = Field(description="Total number of documents for this track_id")
    status_summary: Dict[str, int] = Field(description="Count of documents by status")

    class Config:
        json_schema_extra = {
            "example": {
                "track_id": "upload_20250729_170612_abc123",
                "documents": [
                    {
                        "id": "doc_123456",
                        "content_summary": "Research paper on machine learning",
                        "content_length": 15240,
                        "status": "PROCESSED",
                        "created_at": "2025-03-31T12:34:56",
                        "updated_at": "2025-03-31T12:35:30",
                        "track_id": "upload_20250729_170612_abc123",
                        "chunks_count": 12,
                        "error": None,
                        "metadata": {"author": "John Doe", "year": 2025},
                        "file_path": "research_paper.pdf",
                    }
                ],
                "total_count": 1,
                "status_summary": {"PROCESSED": 1},
            }
        }


class DocumentsRequest(BaseModel):
    """Request model for paginated document queries

    Attributes:
        status_filter: Filter by document status, None for all statuses
        page: Page number (1-based)
        page_size: Number of documents per page (10-200)
        sort_field: Field to sort by ('created_at', 'updated_at', 'id', 'file_path')
        sort_direction: Sort direction ('asc' or 'desc')
    """

    status_filter: Optional[DocStatus] = Field(
        default=None, description="Filter by document status, None for all statuses"
    )
    page: int = Field(default=1, ge=1, description="Page number (1-based)")
    page_size: int = Field(
        default=50, ge=10, le=200, description="Number of documents per page (10-200)"
    )
    sort_field: Literal["created_at", "updated_at", "id", "file_path"] = Field(
        default="updated_at", description="Field to sort by"
    )
    sort_direction: Literal["asc", "desc"] = Field(
        default="desc", description="Sort direction"
    )
    # standard_type: str = Field(default=None, description="可选的标准类型过滤，不填则返回所有")
    class Config:
        json_schema_extra = {
            "example": {
                "status_filter": "PROCESSED",
                "page": 1,
                "page_size": 50,
                "sort_field": "updated_at",
                "sort_direction": "desc",
            }
        }


class PaginationInfo(BaseModel):
    """Pagination information

    Attributes:
        page: Current page number
        page_size: Number of items per page
        total_count: Total number of items
        total_pages: Total number of pages
        has_next: Whether there is a next page
        has_prev: Whether there is a previous page
    """

    page: int = Field(description="Current page number")
    page_size: int = Field(description="Number of items per page")
    total_count: int = Field(description="Total number of items")
    total_pages: int = Field(description="Total number of pages")
    has_next: bool = Field(description="Whether there is a next page")
    has_prev: bool = Field(description="Whether there is a previous page")

    class Config:
        json_schema_extra = {
            "example": {
                "page": 1,
                "page_size": 50,
                "total_count": 150,
                "total_pages": 3,
                "has_next": True,
                "has_prev": False,
            }
        }


class PaginatedDocsResponse(BaseModel):
    """Response model for paginated document queries

    Attributes:
        documents: List of documents for the current page
        pagination: Pagination information
        status_counts: Count of documents by status for all documents
    """

    documents: List[DocStatusResponse] = Field(
        description="List of documents for the current page"
    )
    pagination: PaginationInfo = Field(description="Pagination information")
    status_counts: Dict[str, int] = Field(
        description="Count of documents by status for all documents"
    )

    class Config:
        json_schema_extra = {
            "example": {
                "documents": [
                    {
                        "id": "doc_123456",
                        "content_summary": "Research paper on machine learning",
                        "content_length": 15240,
                        "status": "PROCESSED",
                        "created_at": "2025-03-31T12:34:56",
                        "updated_at": "2025-03-31T12:35:30",
                        "track_id": "upload_20250729_170612_abc123",
                        "chunks_count": 12,
                        "error_msg": None,
                        "metadata": {"author": "John Doe", "year": 2025},
                        "file_path": "research_paper.pdf",
                    }
                ],
                "pagination": {
                    "page": 1,
                    "page_size": 50,
                    "total_count": 150,
                    "total_pages": 3,
                    "has_next": True,
                    "has_prev": False,
                },
                "status_counts": {
                    "PENDING": 10,
                    "PROCESSING": 5,
                    "PREPROCESSED": 5,
                    "PROCESSED": 130,
                    "FAILED": 5,
                },
            }
        }


class StatusCountsResponse(BaseModel):
    """Response model for document status counts

    Attributes:
        status_counts: Count of documents by status
    """

    status_counts: Dict[str, int] = Field(description="Count of documents by status")

    class Config:
        json_schema_extra = {
            "example": {
                "status_counts": {
                    "PENDING": 10,
                    "PROCESSING": 5,
                    "PREPROCESSED": 5,
                    "PROCESSED": 130,
                    "FAILED": 5,
                }
            }
        }


class PipelineStatusResponse(BaseModel):
    """Response model for pipeline status

    Attributes:
        autoscanned: Whether auto-scan has started
        busy: Whether the pipeline is currently busy
        job_name: Current job name (e.g., indexing files/indexing texts)
        job_start: Job start time as ISO format string with timezone (optional)
        docs: Total number of documents to be indexed
        batchs: Number of batches for processing documents
        cur_batch: Current processing batch
        request_pending: Flag for pending request for processing
        latest_message: Latest message from pipeline processing
        history_messages: List of history messages
        update_status: Status of update flags for all namespaces
    """

    autoscanned: bool = False
    busy: bool = False
    job_name: str = "Default Job"
    job_start: Optional[str] = None
    docs: int = 0
    batchs: int = 0
    cur_batch: int = 0
    request_pending: bool = False
    latest_message: str = ""
    history_messages: Optional[List[str]] = None
    update_status: Optional[dict] = None

    @field_validator("job_start", mode="before")
    @classmethod
    def parse_job_start(cls, value):
        """Process datetime and return as ISO format string with timezone"""
        return format_datetime(value)

    class Config:
        extra = "allow"  # Allow additional fields from the pipeline status


class DocumentManager:
    def __init__(
        self,
        input_dir: str,
        workspace: str = "",  # New parameter for workspace isolation
        supported_extensions: tuple = (
            ".txt",
            ".md",
            ".mdx",  # MDX (Markdown + JSX)
            ".pdf",
            ".docx",
            ".pptx",
            ".xlsx",
            ".rtf",  # Rich Text Format
            ".odt",  # OpenDocument Text
            ".tex",  # LaTeX
            ".epub",  # Electronic Publication
            ".html",  # HyperText Markup Language
            ".htm",  # HyperText Markup Language
            ".csv",  # Comma-Separated Values
            ".json",  # JavaScript Object Notation
            ".xml",  # eXtensible Markup Language
            ".yaml",  # YAML Ain't Markup Language
            ".yml",  # YAML
            ".log",  # Log files
            ".conf",  # Configuration files
            ".ini",  # Initialization files
            ".properties",  # Java properties files
            ".sql",  # SQL scripts
            ".bat",  # Batch files
            ".sh",  # Shell scripts
            ".c",  # C source code
            ".h",  # C header
            ".cpp",  # C++ source code
            ".hpp",  # C++ header
            ".py",  # Python source code
            ".java",  # Java source code
            ".js",  # JavaScript source code
            ".ts",  # TypeScript source code
            ".swift",  # Swift source code
            ".go",  # Go source code
            ".rb",  # Ruby source code
            ".php",  # PHP source code
            ".css",  # Cascading Style Sheets
            ".scss",  # Sassy CSS
            ".less",  # LESS CSS
        ),
    ):
        # Store the base input directory and workspace
        self.base_input_dir = Path(input_dir)
        self.workspace = workspace
        self.supported_extensions = supported_extensions
        self.indexed_files = set()

        # Create workspace-specific input directory
        # If workspace is provided, create a subdirectory for data isolation
        if workspace:
            self.input_dir = self.base_input_dir / workspace
        else:
            self.input_dir = self.base_input_dir

        # Create input directory if it doesn't exist
        self.input_dir.mkdir(parents=True, exist_ok=True)

    def scan_directory_for_new_files(self) -> List[Path]:
        """Scan input directory for new files"""
        new_files = []
        for ext in self.supported_extensions:
            logger.debug(f"Scanning for {ext} files in {self.input_dir}")
            for file_path in self.input_dir.glob(f"*{ext}"):
                if file_path not in self.indexed_files:
                    new_files.append(file_path)
        return new_files

    def mark_as_indexed(self, file_path: Path):
        self.indexed_files.add(file_path)

    def is_supported_file(self, filename: str) -> bool:
        return any(filename.lower().endswith(ext) for ext in self.supported_extensions)


def validate_file_path_security(file_path_str: str, base_dir: Path) -> Optional[Path]:
    """
    Validate file path security to prevent Path Traversal attacks.

    Args:
        file_path_str: The file path string to validate
        base_dir: The base directory that the file must be within

    Returns:
        Path: Safe file path if valid, None if unsafe or invalid
    """
    if not file_path_str or not file_path_str.strip():
        return None

    try:
        # Clean the file path string
        clean_path_str = file_path_str.strip()

        # Check for obvious path traversal patterns before processing
        # This catches both Unix (..) and Windows (..\) style traversals
        if ".." in clean_path_str:
            # Additional check for Windows-style backslash traversal
            if (
                "\\..\\" in clean_path_str
                or clean_path_str.startswith("..\\")
                or clean_path_str.endswith("\\..")
            ):
                # logger.warning(
                #     f"Security violation: Windows path traversal attempt detected - {file_path_str}"
                # )
                return None

        # Normalize path separators (convert backslashes to forward slashes)
        # This helps handle Windows-style paths on Unix systems
        normalized_path = clean_path_str.replace("\\", "/")

        # Create path object and resolve it (handles symlinks and relative paths)
        candidate_path = (base_dir / normalized_path).resolve()
        base_dir_resolved = base_dir.resolve()

        # Check if the resolved path is within the base directory
        if not candidate_path.is_relative_to(base_dir_resolved):
            # logger.warning(
            #     f"Security violation: Path traversal attempt detected - {file_path_str}"
            # )
            return None

        return candidate_path

    except (OSError, ValueError, Exception) as e:
        logger.warning(f"Invalid file path detected: {file_path_str} - {str(e)}")
        return None


def get_unique_filename_in_enqueued(target_dir: Path, original_name: str) -> str:
    """Generate a unique filename in the target directory by adding numeric suffixes if needed

    Args:
        target_dir: Target directory path
        original_name: Original filename

    Returns:
        str: Unique filename (may have numeric suffix added)
    """
    import time

    original_path = Path(original_name)
    base_name = original_path.stem
    extension = original_path.suffix

    # Try original name first
    if not (target_dir / original_name).exists():
        return original_name

    # Try with numeric suffixes 001-999
    for i in range(1, 1000):
        suffix = f"{i:03d}"
        new_name = f"{base_name}_{suffix}{extension}"
        if not (target_dir / new_name).exists():
            return new_name

    # Fallback with timestamp if all 999 slots are taken
    timestamp = int(time.time())
    return f"{base_name}_{timestamp}{extension}"


# Document processing helper functions (synchronous)
# These functions run in thread pool via asyncio.to_thread() to avoid blocking the event loop


def _convert_with_docling(file_path: Path) -> str:
    """Convert document using docling (synchronous).

    Args:
        file_path: Path to the document file

    Returns:
        str: Extracted markdown content
    """
    from docling.document_converter import DocumentConverter  # type: ignore

    converter = DocumentConverter()
    result = converter.convert(file_path)
    return result.document.export_to_markdown()


def _extract_table_keywords(text: str, limit: int = 12) -> list[str]:
    keywords: list[str] = []
    seen = set()
    for token in re.findall(r"[\u4e00-\u9fffA-Za-z0-9]{1,12}", text):
        if token in seen:
            continue
        seen.add(token)
        keywords.append(token)
        if len(keywords) >= limit:
            break
    return keywords


def _infer_table_title(lines: list[str], start_idx: int) -> str:
    for i in range(start_idx - 1, -1, -1):
        line = lines[i].strip()
        if not line:
            continue
        if "表" in line:
            return line
        break
    return "TABLE"


def _enhance_markdown_tables(text: str) -> str:
    lines = text.splitlines()
    output: list[str] = []
    i = 0
    while i < len(lines):
        line = lines[i]
        if "|" in line and i + 1 < len(lines):
            sep = lines[i + 1]
            if re.match(r"^\s*\|?(?:\s*:?-+:?\s*\|)+\s*:?-+:?\s*\|?\s*$", sep):
                header = line
                title = _infer_table_title(lines, i)
                header_text = header.replace("|", " ")
                keywords = _extract_table_keywords(header_text)
                tag_line = "[TABLE_TAG] " + title + " " + " ".join(keywords)
                output.append(tag_line.strip())
                output.append(line)
                output.append(sep)
                i += 2
                while i < len(lines) and lines[i].strip():
                    output.append(lines[i])
                    i += 1
                continue
        output.append(line)
        i += 1
    return "\n".join(output)


def _enhance_html_tables(text: str) -> str:
    def _replace(match: re.Match) -> str:
        table_html = match.group(0)
        header_cells = re.findall(r"<th[^>]*>(.*?)</th>", table_html, flags=re.S)
        if not header_cells:
            header_cells = re.findall(r"<td[^>]*>(.*?)</td>", table_html, flags=re.S)
        header_text = re.sub(r"<[^>]+>", " ", " ".join(header_cells))
        keywords = _extract_table_keywords(header_text)
        tag_line = "[TABLE_TAG] TABLE " + " ".join(keywords)
        return tag_line.strip() + "\n" + table_html

    return re.sub(r"<table[\\s\\S]*?</table>", _replace, text, flags=re.S)


def _enhance_tables(text: str) -> str:
    if "<table" in text:
        text = _enhance_html_tables(text)
    if "|" in text:
        text = _enhance_markdown_tables(text)
    return text


def _extract_pdf_pypdf(file_bytes: bytes, password: Optional[str] = None) -> str:
    """Extract PDF content using pypdf (synchronous).

    Args:
        file_bytes: PDF file content as bytes
        password: Optional password for encrypted PDFs

    Returns:
        str: Extracted text content

    Raises:
        Exception: If PDF is encrypted and password is incorrect or missing
    """
    from pypdf import PdfReader  # type: ignore

    pdf_file = BytesIO(file_bytes)
    reader = PdfReader(pdf_file)

    # Check if PDF is encrypted
    if reader.is_encrypted:
        if not password:
            raise Exception("PDF is encrypted but no password provided")

        decrypt_result = reader.decrypt(password)
        if decrypt_result == 0:
            raise Exception("Incorrect PDF password")

    # Extract text from all pages
    content = ""
    for page in reader.pages:
        content += page.extract_text() + "\n"

    return content


def _extract_docx(file_bytes: bytes) -> str:
    """Extract DOCX content including tables in document order (synchronous).

    Args:
        file_bytes: DOCX file content as bytes

    Returns:
        str: Extracted text content with tables in their original positions.
             Tables are separated from paragraphs with blank lines for clarity.
    """
    from docx import Document  # type: ignore
    from docx.table import Table  # type: ignore
    from docx.text.paragraph import Paragraph  # type: ignore

    docx_file = BytesIO(file_bytes)
    doc = Document(docx_file)

    def escape_cell(cell_value: str | None) -> str:
        """Escape characters that would break tab-delimited layout.

        Escape order is critical: backslashes first, then tabs/newlines.
        This prevents double-escaping issues.

        Args:
            cell_value: The cell value to escape (can be None or str)

        Returns:
            str: Escaped cell value safe for tab-delimited format
        """
        if cell_value is None:
            return ""
        text = str(cell_value)
        # CRITICAL: Escape backslash first to avoid double-escaping
        return (
            text.replace("\\", "\\\\")  # Must be first: \ -> \\
            .replace("\t", "&emsp;&emsp;")  # Tab -> \t (visible)
            .replace("\r\n", "<br>")  # Windows newline -> \n
            .replace("\r", "<br>")  # Mac newline -> \n
            .replace("\n", "<br>")  # Unix newline -> \n
        )

    content_parts = []
    in_table = False  # Track if we're currently processing a table

    # Iterate through all body elements in document order
    for element in doc.element.body:
        # Check if element is a paragraph
        if element.tag.endswith("p"):
            # If coming out of a table, add blank line after table
            if in_table:
                content_parts.append("")  # Blank line after table
                in_table = False

            paragraph = Paragraph(element, doc)
            text = paragraph.text
            # Always append to preserve document spacing (including blank paragraphs)
            content_parts.append(text)

        # Check if element is a table
        elif element.tag.endswith("tbl"):
            # Add blank line before table (if content exists)
            if content_parts and not in_table:
                content_parts.append("")  # Blank line before table

            in_table = True
            table = Table(element, doc)
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text
                    # Escape special characters to preserve tab-delimited structure
                    row_text.append(escape_cell(cell_text))
                # Only add row if at least one cell has content
                if any(cell for cell in row_text):
                    content_parts.append("\t".join(row_text))

    return "\n".join(content_parts)


def _extract_pptx(file_bytes: bytes) -> str:
    """Extract PPTX content (synchronous).

    Args:
        file_bytes: PPTX file content as bytes

    Returns:
        str: Extracted text content
    """
    from pptx import Presentation  # type: ignore

    pptx_file = BytesIO(file_bytes)
    prs = Presentation(pptx_file)
    content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
    return content


def _extract_xlsx(file_bytes: bytes) -> str:
    """Extract XLSX content in tab-delimited format with clear sheet separation.

    This function processes Excel workbooks and converts them to a structured text format
    suitable for LLM prompts and RAG systems. Each sheet is clearly delimited with
    separator lines, and special characters are escaped to preserve the tab-delimited structure.

    Features:
    - Each sheet is wrapped with '====================' separators for visual distinction
    - Special characters (tabs, newlines, backslashes) are escaped to prevent structure corruption
    - Column alignment is preserved across all rows to maintain tabular structure
    - Empty rows are preserved as blank lines to maintain row structure
    - Uses sheet.max_column to determine column width efficiently

    Args:
        file_bytes: XLSX file content as bytes

    Returns:
        str: Extracted text content with all sheets in tab-delimited format.
             Format: Sheet separators, sheet name, then tab-delimited rows.

    Example output:
        ==================== Sheet: Data ====================
        Name\tAge\tCity
        Alice\t30\tNew York
        Bob\t25\tLondon

        ==================== Sheet: Summary ====================
        Total\t2
        ====================
    """
    from openpyxl import load_workbook  # type: ignore

    xlsx_file = BytesIO(file_bytes)
    wb = load_workbook(xlsx_file)

    def escape_cell(cell_value: Any) -> str:
        """Escape characters that would break tab-delimited layout.

        Escape order is critical: backslashes first, then tabs/newlines.
        This prevents double-escaping issues.

        Args:
            cell_value: The cell value to escape (can be None, str, int, or float)

        Returns:
            str: Escaped cell value safe for tab-delimited format
        """
        if cell_value is None:
            return ""
        text = str(cell_value)
        # CRITICAL: Escape backslash first to avoid double-escaping
        return (
            text.replace("\\", "\\\\")  # Must be first: \ -> \\
            .replace("\t", "\\t")  # Tab -> \t (visible)
            .replace("\r\n", "\\n")  # Windows newline -> \n
            .replace("\r", "\\n")  # Mac newline -> \n
            .replace("\n", "\\n")  # Unix newline -> \n
        )

    def escape_sheet_title(title: str) -> str:
        """Escape sheet title to prevent formatting issues in separators.

        Args:
            title: Original sheet title

        Returns:
            str: Sanitized sheet title with tabs/newlines replaced
        """
        return str(title).replace("\n", " ").replace("\t", " ").replace("\r", " ")

    content_parts: list[str] = []
    sheet_separator = "=" * 20

    for idx, sheet in enumerate(wb):
        if idx > 0:
            content_parts.append("")  # Blank line between sheets for readability

        # Escape sheet title to handle edge cases with special characters
        safe_title = escape_sheet_title(sheet.title)
        content_parts.append(f"{sheet_separator} Sheet: {safe_title} {sheet_separator}")

        # Use sheet.max_column to get the maximum column width directly
        max_columns = sheet.max_column if sheet.max_column else 0

        # Extract rows with consistent width to preserve column alignment
        for row in sheet.iter_rows(values_only=True):
            row_parts = []

            # Build row up to max_columns width
            for idx in range(max_columns):
                if idx < len(row):
                    row_parts.append(escape_cell(row[idx]))
                else:
                    row_parts.append("")  # Pad short rows

            # Check if row is completely empty
            if all(part == "" for part in row_parts):
                # Preserve empty rows as blank lines (maintains row structure)
                content_parts.append("")
            else:
                # Join all columns to maintain consistent column count
                content_parts.append("\t".join(row_parts))

    # Final separator for symmetry (makes parsing easier)
    content_parts.append(sheet_separator)
    return "\n".join(content_parts)


async def pipeline_enqueue_file(
    rag: LightRAG, file_path: Path, track_id: Optional[str] = None
) -> tuple[bool, str]:
    """Add a file to the queue for processing

    Args:
        rag: LightRAG instance
        file_path: Path to the saved file
        track_id: Optional tracking ID, if not provided will be generated
    Returns:
        tuple: (success: bool, track_id: str)
    """

    # Generate track_id if not provided
    if track_id is None:
        track_id = generate_track_id("unknown")

    try:
        content = ""
        ext = file_path.suffix.lower()
        file_size = 0

        # Get file size for error reporting
        try:
            file_size = file_path.stat().st_size
        except Exception:
            file_size = 0

        file = None
        try:
            async with aiofiles.open(file_path, "rb") as f:
                file = await f.read()
        except PermissionError as e:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "[File Extraction]Permission denied - cannot read file",
                    "original_error": str(e),
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(
                f"[File Extraction]Permission denied reading file: {file_path.name}"
            )
            return False, track_id
        except FileNotFoundError as e:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "[File Extraction]File not found",
                    "original_error": str(e),
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(f"[File Extraction]File not found: {file_path.name}")
            return False, track_id
        except Exception as e:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "[File Extraction]File reading error",
                    "original_error": str(e),
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(
                f"[File Extraction]Error reading file {file_path.name}: {str(e)}"
            )
            return False, track_id

        # Process based on file type
        try:
            match ext:
                case (
                    ".txt"
                    | ".md"
                    | ".mdx"
                    | ".html"
                    | ".htm"
                    | ".tex"
                    | ".json"
                    | ".xml"
                    | ".yaml"
                    | ".yml"
                    | ".rtf"
                    | ".odt"
                    | ".epub"
                    | ".csv"
                    | ".log"
                    | ".conf"
                    | ".ini"
                    | ".properties"
                    | ".sql"
                    | ".bat"
                    | ".sh"
                    | ".c"
                    | ".h"
                    | ".cpp"
                    | ".hpp"
                    | ".py"
                    | ".java"
                    | ".js"
                    | ".ts"
                    | ".swift"
                    | ".go"
                    | ".rb"
                    | ".php"
                    | ".css"
                    | ".scss"
                    | ".less"
                ):
                    try:
                        # Try to decode as UTF-8
                        content = file.decode("utf-8")

                        # Validate content
                        if not content or len(content.strip()) == 0:
                            error_files = [
                                {
                                    "file_path": str(file_path.name),
                                    "error_description": "[File Extraction]Empty file content",
                                    "original_error": "File contains no content or only whitespace",
                                    "file_size": file_size,
                                }
                            ]
                            await rag.apipeline_enqueue_error_documents(
                                error_files, track_id
                            )
                            logger.error(
                                f"[File Extraction]Empty content in file: {file_path.name}"
                            )
                            return False, track_id

                        # Check if content looks like binary data string representation
                        if content.startswith("b'") or content.startswith('b"'):
                            error_files = [
                                {
                                    "file_path": str(file_path.name),
                                    "error_description": "[File Extraction]Binary data in text file",
                                    "original_error": "File appears to contain binary data representation instead of text",
                                    "file_size": file_size,
                                }
                            ]
                            await rag.apipeline_enqueue_error_documents(
                                error_files, track_id
                            )
                            logger.error(
                                f"[File Extraction]File {file_path.name} appears to contain binary data representation instead of text"
                            )
                            return False, track_id

                    except UnicodeDecodeError as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]UTF-8 encoding error, please convert it to UTF-8 before processing",
                                "original_error": f"File is not valid UTF-8 encoded text: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]File {file_path.name} is not valid UTF-8 encoded text. Please convert it to UTF-8 before processing."
                        )
                        return False, track_id

                case ".pdf":
                    try:
                        # Try DOCLING first if configured and available
                        if (
                            global_args.document_loading_engine == "DOCLING"
                            and _is_docling_available()
                        ):
                            content = await asyncio.to_thread(
                                _convert_with_docling, file_path
                            )
                        else:
                            if (
                                global_args.document_loading_engine == "DOCLING"
                                and not _is_docling_available()
                            ):
                                logger.warning(
                                    f"DOCLING engine configured but not available for {file_path.name}. Falling back to pypdf."
                                )
                            # Use pypdf (non-blocking via to_thread)
                            pdf_password = getattr(global_args, "pdf_decrypt_password", None)
                            if not isinstance(pdf_password, str):
                                pdf_password = None
                            content = await asyncio.to_thread(
                                _extract_pdf_pypdf,
                                file,
                                pdf_password,
                            )
                    except Exception as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]PDF processing error",
                                "original_error": f"Failed to extract text from PDF: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]Error processing PDF {file_path.name}: {str(e)}"
                        )
                        return False, track_id

                case ".docx":
                    try:
                        # Try DOCLING first if configured and available
                        if (
                            global_args.document_loading_engine == "DOCLING"
                            and _is_docling_available()
                        ):
                            content = await asyncio.to_thread(
                                _convert_with_docling, file_path
                            )
                        else:
                            if (
                                global_args.document_loading_engine == "DOCLING"
                                and not _is_docling_available()
                            ):
                                logger.warning(
                                    f"DOCLING engine configured but not available for {file_path.name}. Falling back to python-docx."
                                )
                            # Use python-docx (non-blocking via to_thread)
                            content = await asyncio.to_thread(_extract_docx, file)
                    except Exception as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]DOCX processing error",
                                "original_error": f"Failed to extract text from DOCX: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]Error processing DOCX {file_path.name}: {str(e)}"
                        )
                        return False, track_id

                case ".pptx":
                    try:
                        # Try DOCLING first if configured and available
                        if (
                            global_args.document_loading_engine == "DOCLING"
                            and _is_docling_available()
                        ):
                            content = await asyncio.to_thread(
                                _convert_with_docling, file_path
                            )
                        else:
                            if (
                                global_args.document_loading_engine == "DOCLING"
                                and not _is_docling_available()
                            ):
                                logger.warning(
                                    f"DOCLING engine configured but not available for {file_path.name}. Falling back to python-pptx."
                                )
                            # Use python-pptx (non-blocking via to_thread)
                            content = await asyncio.to_thread(_extract_pptx, file)
                    except Exception as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]PPTX processing error",
                                "original_error": f"Failed to extract text from PPTX: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]Error processing PPTX {file_path.name}: {str(e)}"
                        )
                        return False, track_id

                case ".xlsx":
                    try:
                        # Try DOCLING first if configured and available
                        if (
                            global_args.document_loading_engine == "DOCLING"
                            and _is_docling_available()
                        ):
                            content = await asyncio.to_thread(
                                _convert_with_docling, file_path
                            )
                        else:
                            if (
                                global_args.document_loading_engine == "DOCLING"
                                and not _is_docling_available()
                            ):
                                logger.warning(
                                    f"DOCLING engine configured but not available for {file_path.name}. Falling back to openpyxl."
                                )
                            # Use openpyxl (non-blocking via to_thread)
                            content = await asyncio.to_thread(_extract_xlsx, file)
                    except Exception as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]XLSX processing error",
                                "original_error": f"Failed to extract text from XLSX: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]Error processing XLSX {file_path.name}: {str(e)}"
                        )
                        return False, track_id

                case _:
                    error_files = [
                        {
                            "file_path": str(file_path.name),
                            "error_description": f"[File Extraction]Unsupported file type: {ext}",
                            "original_error": f"File extension {ext} is not supported",
                            "file_size": file_size,
                        }
                    ]
                    await rag.apipeline_enqueue_error_documents(error_files, track_id)
                    logger.error(
                        f"[File Extraction]Unsupported file type: {file_path.name} (extension {ext})"
                    )
                    return False, track_id

        except Exception as e:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "[File Extraction]File format processing error",
                    "original_error": f"Unexpected error during file extracting: {str(e)}",
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(
                f"[File Extraction]Unexpected error during {file_path.name} extracting: {str(e)}"
            )
            return False, track_id

        # Insert into the RAG queue
        if content:
            # Check if content contains only whitespace characters
            if not content.strip():
                error_files = [
                    {
                        "file_path": str(file_path.name),
                        "error_description": "[File Extraction]File contains only whitespace",
                        "original_error": "File content contains only whitespace characters",
                        "file_size": file_size,
                    }
                ]
                await rag.apipeline_enqueue_error_documents(error_files, track_id)
                logger.warning(
                    f"[File Extraction]File contains only whitespace characters: {file_path.name}"
                )
                return False, track_id

            try:
                content = _enhance_tables(content)
                await rag.apipeline_enqueue_documents(
                    content, file_paths=file_path.name, track_id=track_id
                )

                logger.info(
                    f"Successfully extracted and enqueued file: {file_path.name}"
                )

                # Move file to __enqueued__ directory after enqueuing
                try:
                    enqueued_dir = file_path.parent / "__enqueued__"
                    enqueued_dir.mkdir(exist_ok=True)

                    # Generate unique filename to avoid conflicts
                    unique_filename = get_unique_filename_in_enqueued(
                        enqueued_dir, file_path.name
                    )
                    target_path = enqueued_dir / unique_filename

                    # Move the file
                    file_path.rename(target_path)
                    logger.debug(
                        f"Moved file to enqueued directory: {file_path.name} -> {unique_filename}"
                    )

                except Exception as move_error:
                    logger.error(
                        f"Failed to move file {file_path.name} to __enqueued__ directory: {move_error}"
                    )
                    # Don't affect the main function's success status

                return True, track_id

            except Exception as e:
                error_files = [
                    {
                        "file_path": str(file_path.name),
                        "error_description": "Document enqueue error",
                        "original_error": f"Failed to enqueue document: {str(e)}",
                        "file_size": file_size,
                    }
                ]
                await rag.apipeline_enqueue_error_documents(error_files, track_id)
                logger.error(f"Error enqueueing document {file_path.name}: {str(e)}")
                return False, track_id
        else:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "No content extracted",
                    "original_error": "No content could be extracted from file",
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(f"No content extracted from file: {file_path.name}")
            return False, track_id

    except Exception as e:
        # Catch-all for any unexpected errors
        try:
            file_size = file_path.stat().st_size if file_path.exists() else 0
        except Exception:
            file_size = 0

        error_files = [
            {
                "file_path": str(file_path.name),
                "error_description": "Unexpected processing error",
                "original_error": f"Unexpected error: {str(e)}",
                "file_size": file_size,
            }
        ]
        await rag.apipeline_enqueue_error_documents(error_files, track_id)
        logger.error(f"Enqueuing file {file_path.name} error: {str(e)}")
        logger.error(traceback.format_exc())
        return False, track_id
    finally:
        if file_path.name.startswith(temp_prefix):
            try:
                file_path.unlink()
            except Exception as e:
                logger.error(f"Error deleting file {file_path}: {str(e)}")


async def pipeline_index_file(rag: LightRAG, file_path: Path, track_id: Optional[str] = None):
    """Index a file with track_id

    Args:
        rag: LightRAG instance
        file_path: Path to the saved file
        track_id: Optional tracking ID
    """
    try:
        success, returned_track_id = await pipeline_enqueue_file(
            rag, file_path, track_id
        )
        if success:
            await rag.apipeline_process_enqueue_documents()

    except Exception as e:
        logger.error(f"Error indexing file {file_path.name}: {str(e)}")
        logger.error(traceback.format_exc())


async def pipeline_index_files(
    rag: LightRAG, file_paths: List[Path], track_id: Optional[str] = None
):
    """Index multiple files sequentially to avoid high CPU load

    Args:
        rag: LightRAG instance
        file_paths: Paths to the files to index
        track_id: Optional tracking ID to pass to all files
    """
    if not file_paths:
        return
    try:
        enqueued = False

        # Use get_pinyin_sort_key for Chinese pinyin sorting
        sorted_file_paths = sorted(
            file_paths, key=lambda p: get_pinyin_sort_key(str(p))
        )

        # Process files sequentially with track_id
        for file_path in sorted_file_paths:
            success, _ = await pipeline_enqueue_file(rag, file_path, track_id)
            if success:
                enqueued = True

        # Process the queue only if at least one file was successfully enqueued
        if enqueued:
            await rag.apipeline_process_enqueue_documents()
    except Exception as e:
        logger.error(f"Error indexing files: {str(e)}")
        logger.error(traceback.format_exc())


async def pipeline_index_texts(
    rag: LightRAG,
    texts: List[str],
    file_sources: Optional[List[str]] = None,
    track_id: Optional[str] = None,
):
    """Index a list of texts with track_id

    Args:
        rag: LightRAG instance
        texts: The texts to index
        file_sources: Sources of the texts
        track_id: Optional tracking ID
    """
    if not texts:
        return
    if file_sources is not None:
        if len(file_sources) != 0 and len(file_sources) != len(texts):
            [
                file_sources.append("unknown_source")
                for _ in range(len(file_sources), len(texts))
            ]
    await rag.apipeline_enqueue_documents(
        input=texts, file_paths=file_sources, track_id=track_id
    )
    await rag.apipeline_process_enqueue_documents()


# async def run_scanning_process(
#     rag_instances: dict, doc_manager: DocumentManager, track_id: str = None
# ):
    # """Background task to scan and index documents

    # Args:
    #     rag: LightRAG instance
    #     doc_manager: DocumentManager instance
    #     track_id: Optional tracking ID to pass to all scanned files
    # """
    # try:
    #     new_files = doc_manager.scan_directory_for_new_files()
    #     if not new_files:
    #         logger.info("No new files to index.")
    #         return 
    #     total_files = len(new_files)
    #     logger.info(f"Found {total_files} files to index.")

    #     if new_files:
    #         # Check for files with PROCESSED status and filter them out
    #         valid_files = []
    #         processed_files = []
    #     # 按标准类型分组，默认分组为 "others"
    #     files_by_type = {"others": []}
        
    #         for file_path in new_files:
    #             filename = file_path.name
    #             existing_doc_data = await rag.doc_status.get_doc_by_file_path(filename)

    #             if existing_doc_data and existing_doc_data.get("status") == "processed":
    #                 # File is already PROCESSED, skip it with warning
    #                 processed_files.append(filename)
    #                 logger.warning(f"Skipping already processed file: {filename}")
    #             else:
    #                 # File is new or in non-PROCESSED status, add to processing list
    #                 valid_files.append(file_path)

    #         # Process valid files (new files + non-PROCESSED status files)
    #         if valid_files:
    #             await pipeline_index_files(rag, valid_files, track_id)
    #             if processed_files:
    #                 logger.info(
    #                     f"Scanning process completed: {len(valid_files)} files Processed {len(processed_files)} skipped."
    #                 )
    #             else:
    #                 logger.info(
    #                     f"Scanning process completed: {len(valid_files)} files Processed."
    #                 )
    #         else:
    #             logger.info(
    #                 "No files to process after filtering already processed files."
    #             )
    #     else:
    #         # No new files to index, check if there are any documents in the queue
    #         logger.info(
    #             "No upload file found, check if there are any documents in the queue..."
    #         )
    #         await rag.apipeline_process_enqueue_documents()

    # except Exception as e:
    #     logger.error(f"Error during scanning process: {str(e)}")
    #     logger.error(traceback.format_exc())

async def run_scanning_process(
    rag_instances: dict,
    doc_manager: DocumentManager,
    track_id: Optional[str] = None
):
    try:
        new_files = doc_manager.scan_directory_for_new_files()
        if not new_files:
            logger.info("No new files to index.")
            return

        # 按标准类型分组，默认分组为 "others"
        files_by_type = {"others": []}
        for file_path in new_files:
            std_type = detect_standard_type(file_path.name)
            if std_type not in rag_instances:
                std_type = "others"
            files_by_type.setdefault(std_type, []).append(file_path)

        # 为每个类型分别启动索引任务
        for std_type, file_list in files_by_type.items():
            if file_list:
                rag = rag_instances[std_type]
                await pipeline_index_files(rag, file_list, track_id)

        logger.info("Scanning process completed.")
    except Exception as e:
        logger.error(f"Error during scanning: {e}")


async def background_delete_documents(
    rag: LightRAG,
    doc_manager: DocumentManager,
    doc_ids: List[str],
    delete_file: bool = False,
    delete_llm_cache: bool = False,
):
    """Background task to delete multiple documents"""
    from lightrag.kg.shared_storage import (
        get_namespace_data,
        get_namespace_lock,
    )

    pipeline_status = await get_namespace_data(
        "pipeline_status", workspace=rag.workspace
    )
    pipeline_status_lock = get_namespace_lock(
        "pipeline_status", workspace=rag.workspace
    )

    total_docs = len(doc_ids)
    successful_deletions = []
    failed_deletions = []

    # Double-check pipeline status before proceeding
    async with pipeline_status_lock:
        if pipeline_status.get("busy", False):
            logger.warning("Error: Unexpected pipeline busy state, aborting deletion.")
            return  # Abort deletion operation

        # Set pipeline status to busy for deletion
        pipeline_status.update(
            {
                "busy": True,
                # Job name can not be changed, it's verified in adelete_by_doc_id()
                "job_name": f"Deleting {total_docs} Documents",
                "job_start": datetime.now().isoformat(),
                "docs": total_docs,
                "batchs": total_docs,
                "cur_batch": 0,
                "latest_message": "Starting document deletion process",
            }
        )
        # Use slice assignment to clear the list in place
        pipeline_status["history_messages"][:] = ["Starting document deletion process"]
        if delete_llm_cache:
            pipeline_status["history_messages"].append(
                "LLM cache cleanup requested for this deletion job"
            )

    try:
        # Loop through each document ID and delete them one by one
        for i, doc_id in enumerate(doc_ids, 1):
            # Check for cancellation at the start of each document deletion
            async with pipeline_status_lock:
                if pipeline_status.get("cancellation_requested", False):
                    cancel_msg = f"Deletion cancelled by user at document {i}/{total_docs}. {len(successful_deletions)} deleted, {total_docs - i + 1} remaining."
                    logger.info(cancel_msg)
                    pipeline_status["latest_message"] = cancel_msg
                    pipeline_status["history_messages"].append(cancel_msg)
                    # Add remaining documents to failed list with cancellation reason
                    failed_deletions.extend(
                        doc_ids[i - 1 :]
                    )  # i-1 because enumerate starts at 1
                    break  # Exit the loop, remaining documents unchanged

                start_msg = f"Deleting document {i}/{total_docs}: {doc_id}"
                logger.info(start_msg)
                pipeline_status["cur_batch"] = i
                pipeline_status["latest_message"] = start_msg
                pipeline_status["history_messages"].append(start_msg)

            file_path = "#"
            try:
                result = await rag.adelete_by_doc_id(
                    doc_id, delete_llm_cache=delete_llm_cache
                )
                file_path = (
                    getattr(result, "file_path", "-") if "result" in locals() else "-"
                )
                if result.status == "success":
                    successful_deletions.append(doc_id)
                    success_msg = (
                        f"Document deleted {i}/{total_docs}: {doc_id}[{file_path}]"
                    )
                    logger.info(success_msg)
                    async with pipeline_status_lock:
                        pipeline_status["history_messages"].append(success_msg)

                    # Handle file deletion if requested and file_path is available
                    if (
                        delete_file
                        and result.file_path
                        and result.file_path != "unknown_source"
                    ):
                        try:
                            deleted_files = []
                            # SECURITY FIX: Use secure path validation to prevent arbitrary file deletion
                            safe_file_path = validate_file_path_security(
                                result.file_path, doc_manager.input_dir
                            )

                            if safe_file_path is None:
                                # Security violation detected - log and skip file deletion
                                security_msg = f"Security violation: Unsafe file path detected for deletion - {result.file_path}"
                                logger.warning(security_msg)
                                async with pipeline_status_lock:
                                    pipeline_status["latest_message"] = security_msg
                                    pipeline_status["history_messages"].append(
                                        security_msg
                                    )
                            else:
                                # check and delete files from input_dir directory
                                if safe_file_path.exists():
                                    try:
                                        safe_file_path.unlink()
                                        deleted_files.append(safe_file_path.name)
                                        file_delete_msg = f"Successfully deleted input_dir file: {result.file_path}"
                                        logger.info(file_delete_msg)
                                        async with pipeline_status_lock:
                                            pipeline_status["latest_message"] = (
                                                file_delete_msg
                                            )
                                            pipeline_status["history_messages"].append(
                                                file_delete_msg
                                            )
                                    except Exception as file_error:
                                        file_error_msg = f"Failed to delete input_dir file {result.file_path}: {str(file_error)}"
                                        logger.debug(file_error_msg)
                                        async with pipeline_status_lock:
                                            pipeline_status["latest_message"] = (
                                                file_error_msg
                                            )
                                            pipeline_status["history_messages"].append(
                                                file_error_msg
                                            )

                                # Also check and delete files from __enqueued__ directory
                                enqueued_dir = doc_manager.input_dir / "__enqueued__"
                                if enqueued_dir.exists():
                                    # SECURITY FIX: Validate that the file path is safe before processing
                                    # Only proceed if the original path validation passed
                                    base_name = Path(result.file_path).stem
                                    extension = Path(result.file_path).suffix

                                    # Search for exact match and files with numeric suffixes
                                    for enqueued_file in enqueued_dir.glob(
                                        f"{base_name}*{extension}"
                                    ):
                                        # Additional security check: ensure enqueued file is within enqueued directory
                                        safe_enqueued_path = (
                                            validate_file_path_security(
                                                enqueued_file.name, enqueued_dir
                                            )
                                        )
                                        if safe_enqueued_path is not None:
                                            try:
                                                enqueued_file.unlink()
                                                deleted_files.append(enqueued_file.name)
                                                logger.info(
                                                    f"Successfully deleted enqueued file: {enqueued_file.name}"
                                                )
                                            except Exception as enqueued_error:
                                                file_error_msg = f"Failed to delete enqueued file {enqueued_file.name}: {str(enqueued_error)}"
                                                logger.debug(file_error_msg)
                                                async with pipeline_status_lock:
                                                    pipeline_status[
                                                        "latest_message"
                                                    ] = file_error_msg
                                                    pipeline_status[
                                                        "history_messages"
                                                    ].append(file_error_msg)
                                        else:
                                            security_msg = f"Security violation: Unsafe enqueued file path detected - {enqueued_file.name}"
                                            logger.warning(security_msg)

                            if deleted_files == []:
                                # Check if the file was supposed to exist but doesn't
                                if safe_file_path and not safe_file_path.exists():
                                    file_info_msg = f"File not found (already deleted or moved): {result.file_path}"
                                    logger.info(file_info_msg)
                                    async with pipeline_status_lock:
                                        pipeline_status["latest_message"] = file_info_msg
                                        pipeline_status["history_messages"].append(
                                            file_info_msg
                                        )
                                else:
                                    file_error_msg = f"File deletion skipped, missing or unsafe file: {result.file_path}"
                                    logger.warning(file_error_msg)
                                    async with pipeline_status_lock:
                                        pipeline_status["latest_message"] = file_error_msg
                                        pipeline_status["history_messages"].append(
                                            file_error_msg
                                        )

                        except Exception as file_error:
                            file_error_msg = f"Failed to delete file {result.file_path}: {str(file_error)}"
                            logger.error(file_error_msg)
                            async with pipeline_status_lock:
                                pipeline_status["latest_message"] = file_error_msg
                                pipeline_status["history_messages"].append(
                                    file_error_msg
                                )
                    elif delete_file:
                        no_file_msg = (
                            f"File deletion skipped, missing file path: {doc_id}"
                        )
                        logger.warning(no_file_msg)
                        async with pipeline_status_lock:
                            pipeline_status["latest_message"] = no_file_msg
                            pipeline_status["history_messages"].append(no_file_msg)
                else:
                    failed_deletions.append(doc_id)
                    error_msg = f"Failed to delete {i}/{total_docs}: {doc_id}[{file_path}] - {result.message}"
                    logger.error(error_msg)
                    async with pipeline_status_lock:
                        pipeline_status["latest_message"] = error_msg
                        pipeline_status["history_messages"].append(error_msg)

            except Exception as e:
                failed_deletions.append(doc_id)
                error_msg = f"Error deleting document {i}/{total_docs}: {doc_id}[{file_path}] - {str(e)}"
                logger.error(error_msg)
                logger.error(traceback.format_exc())
                async with pipeline_status_lock:
                    pipeline_status["latest_message"] = error_msg
                    pipeline_status["history_messages"].append(error_msg)

    except Exception as e:
        error_msg = f"Critical error during batch deletion: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        async with pipeline_status_lock:
            pipeline_status["history_messages"].append(error_msg)
    finally:
        # Final summary and check for pending requests
        async with pipeline_status_lock:
            pipeline_status["busy"] = False
            pipeline_status["pending_requests"] = False  # Reset pending requests flag
            pipeline_status["cancellation_requested"] = (
                False  # Always reset cancellation flag
            )
            completion_msg = f"Deletion completed: {len(successful_deletions)} successful, {len(failed_deletions)} failed"
            pipeline_status["latest_message"] = completion_msg
            pipeline_status["history_messages"].append(completion_msg)

            # Check if there are pending document indexing requests
            has_pending_request = pipeline_status.get("request_pending", False)

        # If there are pending requests, start document processing pipeline
        if has_pending_request:
            try:
                logger.info(
                    "Processing pending document indexing requests after deletion"
                )
                await rag.apipeline_process_enqueue_documents()
            except Exception as e:
                logger.error(f"Error processing pending documents after deletion: {e}")


def create_document_routes(
     rag_instances: dict, doc_manager: DocumentManager, api_key: Optional[str] = None
):
    # Create combined auth dependency for document routes
    combined_auth = get_combined_auth_dependency(api_key)

    # # ===== 新增：缓存不同标准类型的 LightRAG 实例 =====
    # rag_instances = {rag.workspace: rag}  # 默认实例（workspace可能为空字符串）
    # # =================================================

    @router.post(
        "/scan", response_model=ScanResponse, dependencies=[Depends(combined_auth)]
    )
    async def scan_for_new_documents(background_tasks: BackgroundTasks):
        """
        Trigger the scanning process for new documents.

        This endpoint initiates a background task that scans the input directory for new documents
        and processes them. If a scanning process is already running, it returns a status indicating
        that fact.

        Returns:
            ScanResponse: A response object containing the scanning status and track_id
        """
        # Generate track_id with "scan" prefix for scanning operation
        track_id = generate_track_id("scan")

        # Start the scanning process in the background with track_id
        background_tasks.add_task(run_scanning_process, rag_instances, doc_manager, track_id)
        return ScanResponse(
            status="scanning_started",
            message="Scanning process has been initiated in the background",
            track_id=track_id,
        )
    
    @router.post(
        "/upload", response_model=InsertResponse, dependencies=[Depends(combined_auth)]
    )
    async def upload_to_input_dir(
        background_tasks: BackgroundTasks, 
        file: UploadFile = File(...),
        standard_type: Optional[str] = Form(None)
    ):
        """
        Upload a file to the input directory and index it.

        This API endpoint accepts a file through an HTTP POST request, checks if the
        uploaded file is of a supported type, saves it in the specified input directory,
        indexes it for retrieval, and returns a success status with relevant details.

        **File Size Limit:**
        - Configurable via `MAX_UPLOAD_SIZE` environment variable (default: 100MB)
        - Set to `None` or `0` for unlimited upload size
        - Returns HTTP 413 (Request Entity Too Large) if file exceeds limit

        **Duplicate Detection Behavior:**

        This endpoint handles two types of duplicate scenarios differently:

        1. **Filename Duplicate (Synchronous Detection)**:
           - Detected immediately before file processing
           - Returns `status="duplicated"` with the existing document's track_id
           - Two cases:
             - If filename exists in document storage: returns existing track_id
             - If filename exists in file system only: returns empty track_id ("")

        2. **Content Duplicate (Asynchronous Detection)**:
           - Detected during background processing after content extraction
           - Returns `status="success"` with a new track_id immediately
           - The duplicate is detected later when processing the file content
           - Use `/documents/track_status/{track_id}` to check the final result:
             - Document will have `status="FAILED"`
             - `error_msg` contains "Content already exists. Original doc_id: xxx"
             - `metadata.is_duplicate=true` with reference to original document
             - `metadata.original_doc_id` points to the existing document
             - `metadata.original_track_id` shows the original upload's track_id

        **Why Different Behavior?**
        - Filename check is fast (simple lookup), done synchronously
        - Content extraction is expensive (PDF/DOCX parsing), done asynchronously
        - This design prevents blocking the client during expensive operations

        Args:
            background_tasks: FastAPI BackgroundTasks for async processing
            file (UploadFile): The file to be uploaded. It must have an allowed extension.

        Returns:
            InsertResponse: A response object containing the upload status and a message.
                - status="success": File accepted and queued for processing
                - status="duplicated": Filename already exists (see track_id for existing document)

        Raises:
            HTTPException: If the file type is not supported (400), file too large (413), or other errors occur (500).
        """
        try:
            incoming_filename = file.filename
            if not incoming_filename:
                raise HTTPException(status_code=400, detail="Uploaded file must include a filename")

            # ===== 1. 确定标准类型 =====
            if standard_type and standard_type in rag_instances:
                # 使用用户指定的标准类型
                selected_standard_type = standard_type
            else:
                # 自动检测标准类型
                selected_standard_type = detect_standard_type(incoming_filename)
                if selected_standard_type not in rag_instances:
                    selected_standard_type = "others"   # 回退到 others workspace
            
            rag = rag_instances[selected_standard_type]

            # 生成包含标准类型的 track_id
            track_id = generate_track_id(f"upload_{selected_standard_type}")

            logger.info(f"文件 '{incoming_filename}' 使用标准类型: {selected_standard_type} (用户指定: {standard_type is not None})")
            # ===========================
            # Sanitize filename to prevent Path Traversal attacks
            safe_filename = sanitize_filename(incoming_filename, doc_manager.input_dir)

            max_upload_size_raw = getattr(global_args, "max_upload_size", None)
            max_upload_size: Optional[int] = None
            if max_upload_size_raw is not None and not isinstance(max_upload_size_raw, dict):
                try:
                    max_upload_size = int(max_upload_size_raw)
                except (TypeError, ValueError):
                    max_upload_size = None

            # 文件保存到统一输入目录
            file_path = doc_manager.input_dir / safe_filename
           
            if not doc_manager.is_supported_file(safe_filename):
                raise HTTPException(
                    status_code=400,
                    detail=f"Unsupported file type. Supported types: {doc_manager.supported_extensions}",
                )

            # Check file size limit (if configured)
            if max_upload_size is not None and max_upload_size > 0:
                # Safe access to file size (not available in older Starlette versions)
                file_size = getattr(file, "size", None)

                # Pre-flight size check (only if size is available)
                if file_size is not None:
                    if file_size > max_upload_size:
                        raise HTTPException(
                            status_code=413,
                            detail=f"File too large. Maximum size: {max_upload_size / 1024 / 1024:.1f}MB, uploaded: {file_size / 1024 / 1024:.1f}MB",
                        )
                else:
                    # If size not available, we'll check during streaming
                    logger.debug(
                        f"File size not available in UploadFile for {safe_filename}, will check during streaming"
                    )

            # Check if filename already exists in doc_status storage
            # existing_doc_data = await rag.doc_status.get_doc_by_file_path(safe_filename)
            existing_doc_data = await rag.doc_status.get_doc_by_file_path(safe_filename)
            if existing_doc_data:
                status = existing_doc_data.get("status", "unknown")
                existing_track_id = existing_doc_data.get("track_id") or ""
                return InsertResponse(
                    status="duplicated",
                    message=f"File '{safe_filename}' already exists in document storage (Status: {status}).",
                    track_id=existing_track_id,
                )

            # 检查物理文件是否存在（全局唯一，但不同 workspace 可能冲突，可根据需要决定是否允许）
            if file_path.exists():
                return InsertResponse(
                    status="duplicated",
                    message=f"File '{safe_filename}' already exists in the input directory.",
                    track_id="",
                )

            # Async streaming write with size check
            bytes_written = 0
            chunk_size = 1024 * 1024  # 1MB chunks
            needs_cleanup = False

            async with aiofiles.open(file_path, "wb") as out_file:
                while True:
                    # Read chunk from upload stream
                    chunk = await file.read(chunk_size)
                    if not chunk:
                        break

                    # Check size limit during streaming (if not checked before)
                    if max_upload_size is not None and max_upload_size > 0:
                        bytes_written += len(chunk)
                        if bytes_written > max_upload_size:
                            needs_cleanup = True
                            break

                    # Write chunk to file
                    await out_file.write(chunk)

            # Cleanup after file is closed
            if needs_cleanup:
                size_limit = max_upload_size or 0
                try:
                    file_path.unlink()
                except Exception as cleanup_error:
                    logger.error(
                        f"Error cleaning up oversized file {safe_filename}: {cleanup_error}"
                    )

                raise HTTPException(
                    status_code=413,
                    detail=f"File too large. Maximum size: {size_limit / 1024 / 1024:.1f}MB, uploaded: {bytes_written / 1024 / 1024:.1f}MB",
                )

            # Add to background tasks and get track_id
            background_tasks.add_task(pipeline_index_file, rag, file_path, track_id)

            return InsertResponse(
                status="success",
                message=f"File '{safe_filename}' uploaded successfully to {selected_standard_type} workspace.",
                track_id=track_id,
            )

        except HTTPException:
            # Re-raise HTTP exceptions (400, 413, etc.)
            raise
        except Exception as e:
            logger.error(f"Error /documents/upload: {file.filename}: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/text", response_model=InsertResponse, dependencies=[Depends(combined_auth)]
    )
    async def insert_text(
        request: InsertTextRequest, background_tasks: BackgroundTasks
    ):
        """
        Insert text into the RAG system.

        This endpoint allows you to insert text data into the RAG system for later retrieval
        and use in generating responses.

        Args:
            request (InsertTextRequest): The request body containing the text to be inserted.
            background_tasks: FastAPI BackgroundTasks for async processing

        Returns:
            InsertResponse: A response object containing the status of the operation.

        Raises:
            HTTPException: If an error occurs during text processing (500).
        """
        try:
            rag = rag_instances["others"]
            # Check if file_source already exists in doc_status storage
            if (
                request.file_source
                and request.file_source.strip()
                and request.file_source != "unknown_source"
            ):
                existing_doc_data = await rag.doc_status.get_doc_by_file_path(
                    request.file_source
                )
                if existing_doc_data:
                    # Get document status and track_id from existing document
                    status = existing_doc_data.get("status", "unknown")
                    # Use `or ""` to handle both missing key and None value (e.g., legacy rows without track_id)
                    existing_track_id = existing_doc_data.get("track_id") or ""
                    return InsertResponse(
                        status="duplicated",
                        message=f"File source '{request.file_source}' already exists in document storage (Status: {status}).",
                        track_id=existing_track_id,
                    )

            # Check if content already exists by computing content hash (doc_id)
            sanitized_text = sanitize_text_for_encoding(request.text)
            content_doc_id = compute_mdhash_id(sanitized_text, prefix="doc-")
            existing_doc = await rag.doc_status.get_by_id(content_doc_id)
            if existing_doc:
                # Content already exists, return duplicated with existing track_id
                status = existing_doc.get("status", "unknown")
                existing_track_id = existing_doc.get("track_id") or ""
                return InsertResponse(
                    status="duplicated",
                    message=f"Identical content already exists in document storage (doc_id: {content_doc_id}, Status: {status}).",
                    track_id=existing_track_id,
                )

            # Generate track_id for text insertion
            track_id = generate_track_id("insert")

            background_tasks.add_task(
                pipeline_index_texts,
                rag,
                [request.text],
                file_sources=[request.file_source] if request.file_source is not None else None,
                track_id=track_id,
            )

            return InsertResponse(
                status="success",
                message="Text successfully received. Processing will continue in background.",
                track_id=track_id,
            )
        except Exception as e:
            logger.error(f"Error /documents/text: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/texts",
        response_model=InsertResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def insert_texts(
        request: InsertTextsRequest, background_tasks: BackgroundTasks
    ):
        """
        Insert multiple texts into the RAG system.

        This endpoint allows you to insert multiple text entries into the RAG system
        in a single request.

        Args:
            request (InsertTextsRequest): The request body containing the list of texts.
            background_tasks: FastAPI BackgroundTasks for async processing

        Returns:
            InsertResponse: A response object containing the status of the operation.

        Raises:
            HTTPException: If an error occurs during text processing (500).
        """
        try:
            rag = rag_instances["others"]
            # Check if any file_sources already exist in doc_status storage
            if request.file_sources:
                for file_source in request.file_sources:
                    if (
                        file_source
                        and file_source.strip()
                        and file_source != "unknown_source"
                    ):
                        existing_doc_data = await rag.doc_status.get_doc_by_file_path(
                            file_source
                        )
                        if existing_doc_data:
                            # Get document status and track_id from existing document
                            status = existing_doc_data.get("status", "unknown")
                            # Use `or ""` to handle both missing key and None value (e.g., legacy rows without track_id)
                            existing_track_id = existing_doc_data.get("track_id") or ""
                            return InsertResponse(
                                status="duplicated",
                                message=f"File source '{file_source}' already exists in document storage (Status: {status}).",
                                track_id=existing_track_id,
                            )

            # Check if any content already exists by computing content hash (doc_id)
            for text in request.texts:
                sanitized_text = sanitize_text_for_encoding(text)
                content_doc_id = compute_mdhash_id(sanitized_text, prefix="doc-")
                existing_doc = await rag.doc_status.get_by_id(content_doc_id)
                if existing_doc:
                    # Content already exists, return duplicated with existing track_id
                    status = existing_doc.get("status", "unknown")
                    existing_track_id = existing_doc.get("track_id") or ""
                    return InsertResponse(
                        status="duplicated",
                        message=f"Identical content already exists in document storage (doc_id: {content_doc_id}, Status: {status}).",
                        track_id=existing_track_id,
                    )

            # Generate track_id for texts insertion
            track_id = generate_track_id("insert")

            background_tasks.add_task(
                pipeline_index_texts,
                rag,
                request.texts,
                file_sources=request.file_sources,
                track_id=track_id,
            )

            return InsertResponse(
                status="success",
                message="Texts successfully received. Processing will continue in background.",
                track_id=track_id,
            )
        except Exception as e:
            logger.error(f"Error /documents/texts: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    @router.delete(
        "", response_model=ClearDocumentsResponse, dependencies=[Depends(combined_auth)]
    )
    async def clear_documents(standard_type: Optional[str] = Query(None, description="可选的标准类型，不填则清空所有")):
        """
        Clear all documents from the RAG system.
        If standard_type is provided, clear only that workspace; otherwise clear all workspaces.
        """
        from lightrag.kg.shared_storage import (
            get_namespace_data,
            get_namespace_lock,
        )

        async def clear_single_workspace(ws: str, rag_instance: LightRAG) -> tuple[Literal["success", "partial_success", "busy", "fail"], str]:
            """Clear a single workspace and return (status, message)"""
            try:
                pipeline_status = await get_namespace_data(
                    "pipeline_status", workspace=rag_instance.workspace
                )
                pipeline_status_lock = get_namespace_lock(
                    "pipeline_status", workspace=rag_instance.workspace
                )

                # Check and set status with lock
                async with pipeline_status_lock:
                    if pipeline_status.get("busy", False):
                        return ("busy", f"Workspace {ws} is busy, skipped.")

                    pipeline_status.update(
                        {
                            "busy": True,
                            "job_name": "Clearing Documents",
                            "job_start": datetime.now().isoformat(),
                            "docs": 0,
                            "batchs": 0,
                            "cur_batch": 0,
                            "request_pending": False,
                            "latest_message": f"Starting document clearing process for {ws}",
                        }
                    )
                    del pipeline_status["history_messages"][:]
                    pipeline_status["history_messages"].append(f"Starting document clearing process for {ws}")

                # Use drop method to clear all data for this workspace
                drop_tasks = []
                storages = [
                    rag_instance.text_chunks,
                    rag_instance.full_docs,
                    rag_instance.full_entities,
                    rag_instance.full_relations,
                    rag_instance.entity_chunks,
                    rag_instance.relation_chunks,
                    rag_instance.entities_vdb,
                    rag_instance.relationships_vdb,
                    rag_instance.chunks_vdb,
                    rag_instance.chunk_entity_relation_graph,
                    rag_instance.doc_status,
                ]

                for storage in storages:
                    if storage is not None:
                        drop_tasks.append(storage.drop())

                drop_results = await asyncio.gather(*drop_tasks, return_exceptions=True)

                errors = []
                for i, result in enumerate(drop_results):
                    if isinstance(result, Exception):
                        errors.append(f"{storages[i].__class__.__name__}: {str(result)}")

                # Delete files in input directory for this workspace
                workspace_input_dir = doc_manager.base_input_dir / (ws if ws else "others")
                deleted_files = 0
                if workspace_input_dir.exists():
                    for file_path in workspace_input_dir.glob("*"):
                        if file_path.is_file():
                            try:
                                file_path.unlink()
                                deleted_files += 1
                            except Exception as e:
                                errors.append(f"File deletion error: {str(e)}")

                async with pipeline_status_lock:
                    pipeline_status["busy"] = False

                if errors:
                    return ("partial_success", f"Workspace {ws}: cleared with errors: {', '.join(errors)}")
                else:
                    return ("success", f"Workspace {ws}: cleared successfully, deleted {deleted_files} files.")

            except Exception as e:
                logger.error(f"Error clearing workspace {ws}: {e}", exc_info=True)
                return ("fail", f"Workspace {ws}: {str(e)}")

        try:
            if standard_type is not None:
                # 清除单个 workspace
                if standard_type not in rag_instances:
                    raise HTTPException(status_code=400, detail=f"Invalid standard_type: {standard_type}")
                rag = rag_instances[standard_type]
                status, message = await clear_single_workspace(standard_type, rag)
                return ClearDocumentsResponse(status=status, message=message)
            else:
                # 清除所有 workspace
                tasks = [clear_single_workspace(ws, rag) for ws, rag in rag_instances.items()]
                results = await asyncio.gather(*tasks)
                success_count = 0
                partial_count = 0
                fail_count = 0
                messages = []
                for status, msg in results:
                    messages.append(msg)
                    if status == "success":
                        success_count += 1
                    elif status == "partial_success":
                        partial_count += 1
                    else:
                        fail_count += 1

                overall_status = "success" if fail_count == 0 else "fail"
                if success_count > 0 and fail_count > 0:
                    overall_status = "partial_success"

                summary = f"Cleared {success_count} workspaces successfully, {partial_count} partially, {fail_count} failed."
                messages.insert(0, summary)
                return ClearDocumentsResponse(status=overall_status, message="\n".join(messages))

        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Error in clear_documents: {e}", exc_info=True)
            raise HTTPException(status_code=500, detail=str(e))
    @router.get(
        "/pipeline_status",
        dependencies=[Depends(combined_auth)],
        response_model=Dict[str, PipelineStatusResponse],
    )
    async def get_pipeline_status() -> Dict[str, PipelineStatusResponse]:
        """
        Get the current status of the document indexing pipeline.

        This endpoint returns information about the current state of the document processing pipeline,
        including the processing status, progress information, and history messages.

        Returns:
            PipelineStatusResponse: A response object containing:
                - autoscanned (bool): Whether auto-scan has started
                - busy (bool): Whether the pipeline is currently busy
                - job_name (str): Current job name (e.g., indexing files/indexing texts)
                - job_start (str, optional): Job start time as ISO format string
                - docs (int): Total number of documents to be indexed
                - batchs (int): Number of batches for processing documents
                - cur_batch (int): Current processing batch
                - request_pending (bool): Flag for pending request for processing
                - latest_message (str): Latest message from pipeline processing
                - history_messages (List[str], optional): List of history messages (limited to latest 1000 entries,
                  with truncation message if more than 1000 messages exist)

        Raises:
            HTTPException: If an error occurs while retrieving pipeline status (500)
        """
        try:
            from lightrag.kg.shared_storage import get_namespace_data, get_all_update_flags_status

            results = {}
            for ws, rag in rag_instances.items():
                # 获取 pipeline_status 并转换为普通 dict
                pipeline_status = await get_namespace_data("pipeline_status", workspace=rag.workspace)
                status_dict = dict(pipeline_status)

                # 获取 update_flags 并递归转换 MutableBoolean 为 bool
                update_status = await get_all_update_flags_status(workspace=rag.workspace)

                def convert_mutable_bool(obj):
                    """递归转换所有 MutableBoolean 为 bool"""
                    if hasattr(obj, 'value'):  # MutableBoolean 对象
                        return bool(obj.value)
                    elif isinstance(obj, dict):
                        return {k: convert_mutable_bool(v) for k, v in obj.items()}
                    elif isinstance(obj, list):
                        return [convert_mutable_bool(item) for item in obj]
                    else:
                        return obj

                status_dict["update_status"] = convert_mutable_bool(update_status)

                # 处理 history_messages 列表（如有必要）
                if "history_messages" in status_dict:
                    status_dict["history_messages"] = list(status_dict["history_messages"])

                results[ws] = status_dict

            return cast(Dict[str, PipelineStatusResponse], results)

        except Exception as e:
            logger.error(f"Error getting pipeline status: {str(e)}")
            raise HTTPException(status_code=500, detail=str(e))

    # TODO: Deprecated, use /documents/paginated instead
    @router.get(
        "", response_model=DocsStatusesResponse, dependencies=[Depends(combined_auth)]
    )
    async def documents() -> DocsStatusesResponse:
        """
        Get the status of all documents in the system. This endpoint is deprecated; use /documents/paginated instead.
        To prevent excessive resource consumption, a maximum of 1,000 records is returned.

        This endpoint retrieves the current status of all documents, grouped by their
        processing status (PENDING, PROCESSING, PREPROCESSED, PROCESSED, FAILED). The results are
        limited to 1000 total documents with fair distribution across all statuses.

        Returns:
            DocsStatusesResponse: A response object containing a dictionary where keys are
                                DocStatus values and values are lists of DocStatusResponse
                                objects representing documents in each status category.
                                Maximum 1000 documents total will be returned.

        Raises:
            HTTPException: If an error occurs while retrieving document statuses (500).
        """
        try:
            statuses = (
                DocStatus.PENDING,
                DocStatus.PROCESSING,
                DocStatus.PREPROCESSED,
                DocStatus.PROCESSED,
                DocStatus.FAILED,
            )

            tasks = [
                rag_instance.get_docs_by_status(status)
                for rag_instance in rag_instances.values()
                for status in statuses
            ]
            results: List[Dict[str, DocProcessingStatus]] = await asyncio.gather(*tasks)

            response = DocsStatusesResponse()
            total_documents = 0
            max_documents = 1000

            # Convert results to lists for easier processing
            status_documents: list[tuple[DocStatus, list[tuple[str, DocProcessingStatus]]]] = []
            for status_idx, status in enumerate(statuses):
                docs_list: list[tuple[str, DocProcessingStatus]] = []
                for workspace_idx in range(len(rag_instances)):
                    result = results[workspace_idx * len(statuses) + status_idx]
                    docs_list.extend(result.items())
                status_documents.append((status, docs_list))

            # Fair distribution: round-robin across statuses
            status_indices = [0] * len(
                status_documents
            )  # Track current index for each status
            current_status_idx = 0

            while total_documents < max_documents:
                # Check if we have any documents left to process
                has_remaining = False
                for status_idx, (status, docs_list) in enumerate(status_documents):
                    if status_indices[status_idx] < len(docs_list):
                        has_remaining = True
                        break

                if not has_remaining:
                    break

                # Try to get a document from the current status
                status, docs_list = status_documents[current_status_idx]
                current_index = status_indices[current_status_idx]

                if current_index < len(docs_list):
                    doc_id, doc_status = docs_list[current_index]

                    if status not in response.statuses:
                        response.statuses[status] = []

                    response.statuses[status].append(
                        DocStatusResponse(
                            id=doc_id,
                            content_summary=doc_status.content_summary,
                            content_length=doc_status.content_length,
                            status=doc_status.status,
                            created_at=format_datetime(doc_status.created_at),
                            updated_at=format_datetime(doc_status.updated_at),
                            track_id=doc_status.track_id,
                            chunks_count=doc_status.chunks_count,
                            error_msg=doc_status.error_msg,
                            metadata=doc_status.metadata,
                            file_path=doc_status.file_path,
                        )
                    )

                    status_indices[current_status_idx] += 1
                    total_documents += 1

                # Move to next status (round-robin)
                current_status_idx = (current_status_idx + 1) % len(status_documents)

            return response
        except Exception as e:
            logger.error(f"Error GET /documents: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))

    class DeleteDocByIdResponse(BaseModel):
        """Response model for single document deletion operation."""

        status: Literal["deletion_started", "busy", "not_allowed"] = Field(
            description="Status of the deletion operation"
        )
        message: str = Field(description="Message describing the operation result")
        doc_id: str = Field(description="The ID of the document to delete")


    @router.delete(
    "/delete_document",
    response_model=DeleteDocByIdResponse,
    dependencies=[Depends(combined_auth)],
    summary="Delete a document and all its associated data by its ID.",
)
    async def delete_document(
    delete_request: DeleteDocRequest,
    background_tasks: BackgroundTasks,
) -> DeleteDocByIdResponse:

        doc_ids = delete_request.doc_ids
        std_type = delete_request.standard_type

        try:
            from lightrag.kg.shared_storage import (
                get_namespace_data,
                get_namespace_lock,
            )

            # If user specified a workspace, validate and schedule on that workspace only
            if std_type:
                if std_type not in rag_instances:
                    raise HTTPException(status_code=400, detail=f"Invalid standard_type: {std_type}")

                missing: list[str] = []
                for doc_id in doc_ids:
                    doc = await rag_instances[std_type].doc_status.get_by_id(doc_id)
                    if doc is None:
                        missing.append(doc_id)
                if missing:
                    raise HTTPException(
                        status_code=404,
                        detail=f"Document(s) {', '.join(missing)} not found in workspace {std_type}",
                    )

                rag = rag_instances[std_type]
                pipeline_status = await get_namespace_data(
                    "pipeline_status", workspace=rag.workspace
                )
                pipeline_status_lock = get_namespace_lock(
                    "pipeline_status", workspace=rag.workspace
                )

                async with pipeline_status_lock:
                    if pipeline_status.get("busy", False):
                        return DeleteDocByIdResponse(
                            status="busy",
                            message="Cannot delete documents while pipeline is busy",
                            doc_id=", ".join(doc_ids),
                        )

                background_tasks.add_task(
                    background_delete_documents,
                    rag,
                    doc_manager,
                    doc_ids,
                    delete_request.delete_file,
                    delete_request.delete_llm_cache,
                )

                return DeleteDocByIdResponse(
                    status="deletion_started",
                    message=f"Document deletion for {len(doc_ids)} documents has been initiated in workspace {std_type}. Processing will continue in background.",
                    doc_id=", ".join(doc_ids),
                )

            # No workspace specified: detect per-doc workspace and group
            groups: dict[str, list[str]] = {}
            not_found: list[str] = []
            for doc_id in doc_ids:
                found_ws = None
                for ws, rag in rag_instances.items():
                    doc = await rag.doc_status.get_by_id(doc_id)
                    if doc is not None:
                        found_ws = ws
                        break
                if found_ws:
                    groups.setdefault(found_ws, []).append(doc_id)
                else:
                    not_found.append(doc_id)

            if not groups and not_found:
                raise HTTPException(status_code=404, detail=f"Document(s) {', '.join(not_found)} not found in any workspace")

            scheduled: list[str] = []
            skipped_busy: list[str] = []
            for ws, ids in groups.items():
                rag = rag_instances[ws]
                pipeline_status = await get_namespace_data(
                    "pipeline_status", workspace=rag.workspace
                )
                pipeline_status_lock = get_namespace_lock(
                    "pipeline_status", workspace=rag.workspace
                )

                async with pipeline_status_lock:
                    if pipeline_status.get("busy", False):
                        skipped_busy.append(ws)
                    else:
                        background_tasks.add_task(
                            background_delete_documents,
                            rag,
                            doc_manager,
                            ids,
                            delete_request.delete_file,
                            delete_request.delete_llm_cache,
                        )
                        scheduled.append(ws)

            if scheduled:
                message = f"Deletion initiated for workspaces: {', '.join(scheduled)}."
                if skipped_busy:
                    message += f" Skipped busy workspaces: {', '.join(skipped_busy)}."
                if not_found:
                    message += f" Documents not found: {', '.join(not_found)}."
                return DeleteDocByIdResponse(
                    status="deletion_started",
                    message=message,
                    doc_id=", ".join(doc_ids),
                )
            else:
                # Nothing scheduled (likely all busy)
                if not_found and len(not_found) == len(doc_ids):
                    raise HTTPException(status_code=404, detail=f"Document(s) {', '.join(not_found)} not found in any workspace")
                return DeleteDocByIdResponse(
                    status="busy",
                    message=f"Pipeline busy in target workspaces: {', '.join(skipped_busy)}",
                    doc_id=", ".join(doc_ids),
                )

        except HTTPException:
            raise
        except Exception as e:
            error_msg = f"Error initiating document deletion for {delete_request.doc_ids}: {str(e)}"
            logger.error(error_msg)
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=error_msg)

    @router.post(
        "/clear_cache",
        response_model=ClearCacheResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def clear_cache(request: ClearCacheRequest):
        """
        Clear all cache data from the LLM response cache storage.

        This endpoint clears all cached LLM responses regardless of mode.
        The request body is accepted for API compatibility but is ignored.

        Args:
            request (ClearCacheRequest): The request body (ignored for compatibility).

        Returns:
            ClearCacheResponse: A response object containing the status and message.

        Raises:
            HTTPException: If an error occurs during cache clearing (500).
        """
        try:
            for rag in rag_instances.values():
                await rag.aclear_cache()
            message = "Successfully cleared cache for all workspaces"
            return ClearCacheResponse(status="success", message=message)
        except Exception as e:
            logger.error(f"Error clearing cache: {str(e)}")
            raise HTTPException(status_code=500, detail=str(e))

    @router.delete(
        "/delete_entity",
        response_model=DeletionResult,
        dependencies=[Depends(combined_auth)],
    )
    async def delete_entity(request: DeleteEntityRequest):
        """
        Delete an entity and all its relationships from the knowledge graph.

        Args:
            request (DeleteEntityRequest): The request body containing the entity name.

        Returns:
            DeletionResult: An object containing the outcome of the deletion process.

        Raises:
            HTTPException: If the entity is not found (404) or an error occurs (500).
        """
        try:
            if request.standard_type not in rag_instances:
                raise HTTPException(
                    status_code=400,
                    detail=f"Invalid standard_type: {request.standard_type}",
                )

            rag = rag_instances[request.standard_type]
            exists = await rag.chunk_entity_relation_graph.has_node(request.entity_name)
            if not exists:
                raise HTTPException(
                    status_code=404,
                    detail=f"Entity '{request.entity_name}' not found in workspace {request.standard_type}",
                )

            result = await rag.adelete_by_entity(entity_name=request.entity_name)
            if result.status == "not_found":
                raise HTTPException(status_code=404, detail=result.message)
            if result.status == "fail":
                raise HTTPException(status_code=500, detail=result.message)
            # Set doc_id to empty string since this is an entity operation, not document
            result.doc_id = ""
            return result
        except HTTPException:
            raise
        except Exception as e:
            error_msg = f"Error deleting entity '{request.entity_name}': {str(e)}"
            logger.error(error_msg)
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=error_msg)

    @router.delete(
        "/delete_relation",
        response_model=DeletionResult,
        dependencies=[Depends(combined_auth)],
    )
    async def delete_relation(request: DeleteRelationRequest):
        """
        Delete a relationship between two entities from the knowledge graph.

        Args:
            request (DeleteRelationRequest): The request body containing the source and target entity names.

        Returns:
            DeletionResult: An object containing the outcome of the deletion process.

        Raises:
            HTTPException: If the relation is not found (404) or an error occurs (500).
        """
        try:
            if request.standard_type not in rag_instances:
                raise HTTPException(
                    status_code=400,
                    detail=f"Invalid standard_type: {request.standard_type}",
                )

            rag = rag_instances[request.standard_type]
            exists = await rag.chunk_entity_relation_graph.has_edge(
                request.source_entity,
                request.target_entity,
            )
            if not exists:
                raise HTTPException(
                    status_code=404,
                    detail=(
                        f"Relation from '{request.source_entity}' to '{request.target_entity}' "
                        f"not found in workspace {request.standard_type}"
                    ),
                )

            result = await rag.adelete_by_relation(
                source_entity=request.source_entity,
                target_entity=request.target_entity,
            )
            if result.status == "not_found":
                raise HTTPException(status_code=404, detail=result.message)
            if result.status == "fail":
                raise HTTPException(status_code=500, detail=result.message)
            # Set doc_id to empty string since this is a relation operation, not document
            result.doc_id = ""
            return result
        except HTTPException:
            raise
        except Exception as e:
            error_msg = f"Error deleting relation from '{request.source_entity}' to '{request.target_entity}': {str(e)}"
            logger.error(error_msg)
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=error_msg)

    @router.get(
        "/track_status/{track_id}",
        response_model=TrackStatusResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def get_track_status(track_id: str) -> TrackStatusResponse:
        """
        Get the processing status of documents by tracking ID.

        This endpoint retrieves all documents associated with a specific tracking ID,
        allowing users to monitor the processing progress of their uploaded files or inserted texts.

        Args:
            track_id (str): The tracking ID returned from upload, text, or texts endpoints

        Returns:
            TrackStatusResponse: A response object containing:
                - track_id: The tracking ID
                - documents: List of documents associated with this track_id
                - total_count: Total number of documents for this track_id

        Raises:
            HTTPException: If track_id is invalid (400) or an error occurs (500).
        """
        try:
            parts = track_id.split('_')
            if len(parts) >= 2 and parts[1] in rag_instances:
                std_type = parts[1]
            else:
                std_type = "others"
            rag = rag_instances[std_type]
            # Validate track_id
            if not track_id or not track_id.strip():
                raise HTTPException(status_code=400, detail="Track ID cannot be empty")

            track_id = track_id.strip()

            # Get documents by track_id
            docs_by_track_id = await rag.aget_docs_by_track_id(track_id)

            # Convert to response format
            documents = []
            status_summary = {}

            for doc_id, doc_status in docs_by_track_id.items():
                documents.append(
                    DocStatusResponse(
                        id=doc_id,
                        content_summary=doc_status.content_summary,
                        content_length=doc_status.content_length,
                        status=doc_status.status,
                        created_at=format_datetime(doc_status.created_at),
                        updated_at=format_datetime(doc_status.updated_at),
                        track_id=doc_status.track_id,
                        chunks_count=doc_status.chunks_count,
                        error_msg=doc_status.error_msg,
                        metadata=doc_status.metadata,
                        file_path=doc_status.file_path,
                    )
                )

                # Build status summary
                # Handle both DocStatus enum and string cases for robust deserialization
                status_key = str(doc_status.status)
                status_summary[status_key] = status_summary.get(status_key, 0) + 1

            return TrackStatusResponse(
                track_id=track_id,
                documents=documents,
                total_count=len(documents),
                status_summary=status_summary,
            )

        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Error getting track status for {track_id}: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))
    
    @router.post(
    "/paginated",
    response_model=PaginatedDocsResponse,
    dependencies=[Depends(combined_auth)],
)
    async def get_documents_paginated(request: DocumentsRequest,) -> PaginatedDocsResponse:
        try:
    # 始终聚合所有 workspace
            all_docs = []
            all_status_counts = {}
            for ws, rag in rag_instances.items():
                page = 1
                page_size = 1000
                while True:
                    docs, total = await rag.doc_status.get_docs_paginated(
                        status_filter=request.status_filter,
                        page=page,
                        page_size=page_size,
                        sort_field=request.sort_field,
                        sort_direction=request.sort_direction,
                    )
                    for doc_id, doc in docs:
                        all_docs.append((doc_id, doc, ws))
                    if page * page_size >= total:
                        break
                    page += 1
                counts = await rag.doc_status.get_all_status_counts()
                for status, cnt in counts.items():
                    all_status_counts[status] = all_status_counts.get(status, 0) + cnt

            # 全局排序
            reverse = (request.sort_direction == "desc")
            if request.sort_field == "id":
                all_docs.sort(key=lambda x: x[0], reverse=reverse)
            else:
                all_docs.sort(key=lambda x: getattr(x[1], request.sort_field), reverse=reverse)

            total_count = len(all_docs)
            start = (request.page - 1) * request.page_size
            end = start + request.page_size
            paginated_docs = all_docs[start:end]

            doc_responses = []
            for doc_id, doc, ws in paginated_docs:
                doc_dict = doc.__dict__.copy()
                doc_dict['workspace'] = ws
                doc_responses.append(DocStatusResponse(id=doc_id, **doc_dict))

            total_pages = (total_count + request.page_size - 1) // request.page_size
            pagination = PaginationInfo(
                page=request.page,
                page_size=request.page_size,
                total_count=total_count,
                total_pages=total_pages,
                has_next=request.page < total_pages,
                has_prev=request.page > 1,
            )
            return PaginatedDocsResponse(
                documents=doc_responses,
                pagination=pagination,
                status_counts=all_status_counts,
            )
        except Exception as e:
            logger.error(f"Error getting paginated documents: {str(e)}")
            raise HTTPException(status_code=500, detail=str(e))

    
    @router.get(
        "/status_counts",
        response_model=StatusCountsResponse,
        dependencies=[Depends(combined_auth)],
    )
#     @router.get(
#     "/status_counts",
#     response_model=StatusCountsResponse,
#     dependencies=[Depends(combined_auth)],
# )
    async def get_document_status_counts(
        standard_type: Optional[str] = Query(None, description="可选的标准类型过滤，不填则返回所有")
    ) -> StatusCountsResponse:
        try:
            if standard_type is not None:
                if standard_type not in rag_instances:
                    raise HTTPException(status_code=400, detail=f"Invalid standard_type: {standard_type}")
                rag = rag_instances[standard_type]
                status_counts = await rag.doc_status.get_all_status_counts()
            else:
                # 合并所有 workspace
                status_counts = {}
                for rag in rag_instances.values():
                    counts = await rag.doc_status.get_all_status_counts()
                    for status, cnt in counts.items():
                        status_counts[status] = status_counts.get(status, 0) + cnt
            return StatusCountsResponse(status_counts=status_counts)
        except Exception as e:
            logger.error(f"Error getting document status counts: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=500, detail=str(e))
    # async def get_document_status_counts(standard_type: str = Query("others")) -> StatusCountsResponse:
    #     """
    #     Get counts of documents by status.

    #     This endpoint retrieves the count of documents in each processing status
    #     (PENDING, PROCESSING, PROCESSED, FAILED) for all documents in the system.

    #     Returns:
    #         StatusCountsResponse: A response object containing status counts

    #     Raises:
    #         HTTPException: If an error occurs while retrieving status counts (500).
    #     """
    #     try:
    #         if standard_type not in rag_instances:
    #             raise HTTPException(status_code=400, detail=f"Invalid standard_type: {standard_type}")
    #         rag = rag_instances[standard_type]
    #         status_counts = await rag.doc_status.get_all_status_counts()
    #         return StatusCountsResponse(status_counts=status_counts)

    #     except Exception as e:
    #         logger.error(f"Error getting document status counts: {str(e)}")
    #         logger.error(traceback.format_exc())
    #         raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/reprocess_failed",
        response_model=ReprocessResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def reprocess_failed_documents(background_tasks: BackgroundTasks):
        """
        Reprocess failed and pending documents.

        This endpoint triggers the document processing pipeline which automatically
        picks up and reprocesses documents in the following statuses:
        - FAILED: Documents that failed during previous processing attempts
        - PENDING: Documents waiting to be processed
        - PROCESSING: Documents with abnormally terminated processing (e.g., server crashes)

        This is useful for recovering from server crashes, network errors, LLM service
        outages, or other temporary failures that caused document processing to fail.

        The processing happens in the background and can be monitored by checking the
        pipeline status. The reprocessed documents retain their original track_id from
        initial upload, so use their original track_id to monitor progress.

        Returns:
            ReprocessResponse: Response with status and message.
                track_id is always empty string because reprocessed documents retain
                their original track_id from initial upload.

        Raises:
            HTTPException: If an error occurs while initiating reprocessing (500).
        """
        try:
            for rag in rag_instances.values():
                background_tasks.add_task(rag.apipeline_process_enqueue_documents)
            logger.info("Reprocessing of failed documents initiated for all workspaces")
            return ReprocessResponse(
                status="reprocessing_started",
                message="Reprocessing of failed documents has been initiated in background for all workspaces. Documents retain their original track_id.",
            )
        except Exception as e:
            logger.error(f"Error initiating reprocessing: {str(e)}")
            raise HTTPException(status_code=500, detail=str(e))

    @router.post(
        "/cancel_pipeline",
        response_model=CancelPipelineResponse,
        dependencies=[Depends(combined_auth)],
    )
    async def cancel_pipeline():

        """
        Request cancellation of the currently running pipeline.

        This endpoint sets a cancellation flag in the pipeline status. The pipeline will:
        1. Check this flag at key processing points
        2. Stop processing new documents
        3. Cancel all running document processing tasks
        4. Mark all PROCESSING documents as FAILED with reason "User cancelled"

        The cancellation is graceful and ensures data consistency. Documents that have
        completed processing will remain in PROCESSED status.

        Returns:
            CancelPipelineResponse: Response with status and message
                - status="cancellation_requested": Cancellation flag has been set
                - status="not_busy": Pipeline is not currently running

        Raises:
            HTTPException: If an error occurs while setting cancellation flag (500).
        """
        try:
            from lightrag.kg.shared_storage import get_namespace_data, get_namespace_lock
            responses = []
            for ws, rag in rag_instances.items():
                pipeline_status = await get_namespace_data("pipeline_status", workspace=rag.workspace)
                pipeline_status_lock = get_namespace_lock("pipeline_status", workspace=rag.workspace)
                async with pipeline_status_lock:
                    if pipeline_status.get("busy", False):
                        pipeline_status["cancellation_requested"] = True
                        responses.append(f"Cancellation requested for workspace {ws}")
                    else:
                        responses.append(f"Workspace {ws} not busy")
            return CancelPipelineResponse(
                status="cancellation_requested" if any("requested" in r for r in responses) else "not_busy",
                message="; ".join(responses)
            )
        except Exception as e:
            logger.error(f"Error requesting pipeline cancellation: {str(e)}")
            raise HTTPException(status_code=500, detail=str(e))

    return router
